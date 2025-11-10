"""
Document Processing System for Sales Knowledge
Parses documents (PDFs, Word, Excel, text files) and stores them in vector memory
"""

import os
import json
from typing import List, Dict, Any, Optional
from datetime import datetime
import hashlib

# Document parsing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    import cv2
    import numpy as np
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from Backend.SalesMemory import sales_memory_manager, learn_from_docs

def process_text_file(file_path: str, source_name: Optional[str] = None) -> Dict[str, Any]:
    """
    Process a text file and extract content
    
    Args:
        file_path: Path to the text file
        source_name: Optional custom name for the source
        
    Returns:
        Dictionary with processing results
    """
    try:
        # Try UTF-8 first
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except UnicodeDecodeError:
        # Fallback to latin-1
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
        except Exception as e:
            return {
                "success": False,
                "error": f"Error reading text file: {e}",
                "entries_created": 0
            }
    
    if not source_name:
        source_name = os.path.basename(file_path)
    
    # Split into chunks (to avoid huge single entries)
    chunk_size = 1000  # characters per chunk
    chunks = [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]
    
    entry_ids = []
    for i, chunk in enumerate(chunks):
        if chunk.strip():
            entry_id = learn_from_docs(chunk, f"{source_name}_chunk_{i+1}", "document")
            entry_ids.append(entry_id)
    
    return {
        "success": True,
        "source": source_name,
        "entries_created": len(entry_ids),
        "entry_ids": entry_ids
    }

def process_pdf(file_path: str, source_name: Optional[str] = None, max_pages: int = 50) -> Dict[str, Any]:
    """
    Process a PDF file and extract text content
    
    Args:
        file_path: Path to the PDF file
        source_name: Optional custom name for the source
        max_pages: Maximum number of pages to process
        
    Returns:
        Dictionary with processing results
    """
    if not PDF_AVAILABLE:
        return {
            "success": False,
            "error": "PDF processing not available. Install PyPDF2.",
            "entries_created": 0
        }
    
    try:
        if not source_name:
            source_name = os.path.basename(file_path)
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            total_pages = min(len(pdf_reader.pages), max_pages)
            
            entry_ids = []
            for page_num in range(total_pages):
                try:
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    
                    if text.strip():
                        entry_id = learn_from_docs(
                            text,
                            f"{source_name}_page_{page_num+1}",
                            "document"
                        )
                        entry_ids.append(entry_id)
                except Exception as e:
                    print(f"Error processing page {page_num+1}: {e}")
                    continue
        
        return {
            "success": True,
            "source": source_name,
            "entries_created": len(entry_ids),
            "total_pages": total_pages,
            "entry_ids": entry_ids
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": f"Error processing PDF: {e}",
            "entries_created": 0
        }

def process_word_document(file_path: str, source_name: Optional[str] = None) -> Dict[str, Any]:
    """
    Process a Word document (.docx) and extract text
    
    Args:
        file_path: Path to the Word document
        source_name: Optional custom name for the source
        
    Returns:
        Dictionary with processing results
    """
    if not DOCX_AVAILABLE:
        return {
            "success": False,
            "error": "Word document processing not available. Install python-docx.",
            "entries_created": 0
        }
    
    try:
        if not source_name:
            source_name = os.path.basename(file_path)
        
        doc = Document(file_path)
        
        # Extract text from paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        content = "\n\n".join(paragraphs)
        
        # Split into chunks if large
        chunk_size = 2000
        chunks = [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]
        
        entry_ids = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():
                entry_id = learn_from_docs(
                    chunk,
                    f"{source_name}_section_{i+1}",
                    "document"
                )
                entry_ids.append(entry_id)
        
        return {
            "success": True,
            "source": source_name,
            "entries_created": len(entry_ids),
            "entry_ids": entry_ids
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": f"Error processing Word document: {e}",
            "entries_created": 0
        }

def process_excel_file(file_path: str, source_name: Optional[str] = None, max_sheets: int = 5) -> Dict[str, Any]:
    """
    Process an Excel file and extract data
    
    Args:
        file_path: Path to the Excel file
        source_name: Optional custom name for the source
        max_sheets: Maximum number of sheets to process
        
    Returns:
        Dictionary with processing results
    """
    if not EXCEL_AVAILABLE:
        return {
            "success": False,
            "error": "Excel processing not available. Install pandas and openpyxl.",
            "entries_created": 0
        }
    
    try:
        if not source_name:
            source_name = os.path.basename(file_path)
        
        excel_file = pd.ExcelFile(file_path)
        sheet_names = excel_file.sheet_names[:max_sheets]
        
        entry_ids = []
        for sheet_name in sheet_names:
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to string representation
                # Limit rows and columns for readability
                df_str = df.head(50).to_string(max_cols=20)
                
                content = f"Sheet: {sheet_name}\n\n{df_str}"
                
                entry_id = learn_from_docs(
                    content,
                    f"{source_name}_{sheet_name}",
                    "spreadsheet"
                )
                entry_ids.append(entry_id)
            except Exception as e:
                print(f"Error processing sheet {sheet_name}: {e}")
                continue
        
        return {
            "success": True,
            "source": source_name,
            "entries_created": len(entry_ids),
            "sheets_processed": len(sheet_names),
            "entry_ids": entry_ids
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": f"Error processing Excel file: {e}",
            "entries_created": 0
        }

def process_ppt(file_path: str, source_name: Optional[str] = None) -> Dict[str, Any]:
    """
    Process a PowerPoint file and extract text content
    
    Args:
        file_path: Path to the PPT/PPTX file
        source_name: Optional custom name for the source
        
    Returns:
        Dictionary with processing results
    """
    if not PPTX_AVAILABLE:
        return {
            "success": False,
            "error": "python-pptx library not installed. Install with: pip install python-pptx",
            "entries_created": 0
        }
    
    try:
        # Open the presentation
        prs = Presentation(file_path)
        
        # Extract text from all slides
        all_text = []
        slide_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_count += 1
            slide_text = f"Slide {slide_num}:\n"
            
            # Extract text from shapes (text boxes, placeholders, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
                # Also check for tables
                if hasattr(shape, "table"):
                    table_text = ""
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        table_text += row_text + "\n"
                    if table_text:
                        slide_text += "Table:\n" + table_text
            
            if slide_text.strip() and slide_text.strip() != f"Slide {slide_num}:\n":
                all_text.append(slide_text)
        
        if not all_text:
            return {
                "success": False,
                "error": "No text content found in PowerPoint file",
                "entries_created": 0
            }
        
        # Combine all text
        full_content = "\n\n".join(all_text)
        
        # Store in memory
        source = source_name or os.path.basename(file_path)
        entries_created = learn_from_docs(full_content, source)
        
        return {
            "success": True,
            "source": source,
            "content": full_content,
            "slides_processed": slide_count,
            "entries_created": entries_created
        }
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"Error processing PPT file: {e}")
        print(f"Traceback: {error_trace}")
        return {
            "success": False,
            "error": f"Error processing PowerPoint file: {str(e)}",
            "entries_created": 0
        }

def process_image(file_path: str, source_name: Optional[str] = None) -> Dict[str, Any]:
    """
    Process an image file using OCR to extract text
    
    Args:
        file_path: Path to the image file
        source_name: Optional custom name for the source
        
    Returns:
        Dictionary with processing results
    """
    if not OCR_AVAILABLE:
        return {
            "success": False,
            "error": "Image OCR processing not available. Install pytesseract, PIL, cv2, and numpy.",
            "entries_created": 0
        }
    
    try:
        if not source_name:
            source_name = os.path.basename(file_path)
        
        # Load image
        image = cv2.imread(file_path)
        if image is None:
            return {
                "success": False,
                "error": "Could not read image file",
                "entries_created": 0
            }
        
        # Convert to RGB
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        # Extract text using OCR
        content = pytesseract.image_to_string(image_rgb)
        
        if not content.strip():
            return {
                "success": False,
                "error": "No text detected in image",
                "entries_created": 0
            }
        
        # Split into chunks if large
        chunk_size = 2000
        chunks = [content[i:i+chunk_size] for i in range(0, len(content), chunk_size)]
        
        entry_ids = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():
                entry_id = learn_from_docs(
                    chunk,
                    f"{source_name}_ocr_chunk_{i+1}",
                    "document"
                )
                entry_ids.append(entry_id)
        
        return {
            "success": True,
            "source": source_name,
            "entries_created": len(entry_ids),
            "entry_ids": entry_ids
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": f"Error processing image: {e}",
            "entries_created": 0
        }

def process_document(file_path: str, source_name: Optional[str] = None) -> Dict[str, Any]:
    """
    Process any document type (auto-detects file type)
    
    Args:
        file_path: Path to the document
        source_name: Optional custom name for the source
        
    Returns:
        Dictionary with processing results
    """
    if not os.path.exists(file_path):
        return {
            "success": False,
            "error": f"File not found: {file_path}",
            "entries_created": 0
        }
    
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # If no extension, try to detect file type by reading first bytes
    if not file_ext:
        try:
            with open(file_path, 'rb') as f:
                first_bytes = f.read(4)
                # Check for PDF magic bytes
                if first_bytes.startswith(b'%PDF'):
                    file_ext = '.pdf'
                # Check for ZIP-based formats (docx, xlsx)
                elif first_bytes.startswith(b'PK\x03\x04'):
                    # Try to determine if it's docx or xlsx by checking internal structure
                    f.seek(0)
                    content = f.read(1024)
                    if b'word/' in content:
                        file_ext = '.docx'
                    elif b'xl/' in content or b'worksheets/' in content:
                        file_ext = '.xlsx'
                    else:
                        file_ext = '.zip'
                # Check for text files
                else:
                    try:
                        with open(file_path, 'r', encoding='utf-8') as tf:
                            tf.read(100)  # Try to read as text
                        file_ext = '.txt'
                    except:
                        pass
        except:
            pass
    
    if file_ext == '.pdf':
        return process_pdf(file_path, source_name)
    elif file_ext in ['.docx', '.doc']:
        return process_word_document(file_path, source_name)
    elif file_ext in ['.xlsx', '.xls']:
        return process_excel_file(file_path, source_name)
    elif file_ext in ['.txt', '.md', '.csv', '.log']:
        return process_text_file(file_path, source_name)
    elif file_ext in ['.ppt', '.pptx']:
        return process_ppt(file_path, source_name)
    elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
        return process_image(file_path, source_name)
    else:
        return {
            "success": False,
            "error": f"Unsupported file type: {file_ext}",
            "entries_created": 0
        }


