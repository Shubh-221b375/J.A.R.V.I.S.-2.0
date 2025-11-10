from PyQt5.QtWidgets import (QApplication, QMainWindow, QTextBrowser, QStackedWidget, QWidget, QLineEdit, QGridLayout, QVBoxLayout, QHBoxLayout, QPushButton, QFrame, QLabel, QSizePolicy, QComboBox, QFileDialog, QTextEdit)
from PyQt5.QtGui import QIcon, QPainter, QMovie, QColor, QTextCharFormat, QFont, QPixmap, QTextBlockFormat, QTextOption
from PyQt5.QtCore import Qt, QSize, QTimer
from dotenv import dotenv_values
import sys
import os

# Load environment variables
env_vars = dotenv_values(".env")
Assistantname = env_vars.get("Assistantname")
old_chat_message = ""
# Directory paths
current_dir = os.getcwd()
TempDirPath = rf"{current_dir}\Frontend\Files"
GraphicsDirPath = rf"{current_dir}\Frontend\Graphics"

def AnswerModifier(Answer):
    lines = Answer.split('\n')
    non_empty_lines = [line.strip() for line in lines if line.strip()]
    modified_answer = '\n'.join(non_empty_lines)
    return modified_answer


def QueryModifier(Query):
    new_query = Query.lower().strip()
    query_words  = new_query.split()
    question_words = ['how','what','who','where','when','why','which','whom','can you',"what's", "where's","how's"]

    if any(word + " " in new_query for word in question_words):
        if query_words[-1][-1] in ['.','?','!']:
            new_query = new_query[:-1] + "?"
        else:
            new_query += "?"
    else:
        if query_words[-1][-1] in ['.','?','!']:
            new_query = new_query[:-1] + '.'
        else:
            new_query += '.'

    return new_query.capitalize()


def SetMicrophoneStatus(Command):
    with open(TempDirectoryPath('Mic.data'), 'w', encoding='utf-8') as file:
        file.write(Command)
    

def GetMicrophoneStatus():
    
    with open(TempDirectoryPath('Mic.data'), 'r', encoding='utf-8') as file:
        Status = file.read().strip()
    return Status


def SetAsssistantStatus(Status):
    with open(rf'{TempDirPath}\Status.data','w',encoding='utf-8') as file:
        file.write(Status)


def GetAssistantStatus():
    with open(rf'{TempDirPath}\Status.data', 'r', encoding='utf-8') as file:
        Status = file.read()
    return Status
    

    
# Define placeholders for the missing functions
def MicButtonInitiated():
    SetMicrophoneStatus("True")  # When mic is ON, set status to True
    # Immediately interrupt any ongoing speech when mic is turned on
    try:
        from Backend.TextToSpeech import interrupt_speech
        from Main import force_interrupt_speech
        interrupt_speech()
        force_interrupt_speech()
    except:
        pass

def MicButtonClosed():
    SetMicrophoneStatus("False")  # When mic is OFF, set status to False
    # Also clear any ongoing speech when mic is turned off
    try:
        from Backend.TextToSpeech import interrupt_speech
        interrupt_speech()
    except:
        pass

def MicButtonInitialized():
    """Alias for MicButtonInitiated for compatibility"""
    return MicButtonInitiated()

def GraphicsDirectoryPath(Filename):
    path = rf'{GraphicsDirPath}\{Filename}'
    return path


def TempDirectoryPath(Filename):
    path = rf'{TempDirPath}\{Filename}'
    return path

def ShowTextToScreen(Text):
    with open (rf'{TempDirPath}\Responses.data','w', encoding='utf-8') as file:
        file.write(Text)

    
class ChatSection(QWidget):
    def __init__(self):
        super(ChatSection, self).__init__()
        layout = QVBoxLayout(self)
        # Responsive margins - adjust based on screen size
        layout.setContentsMargins(10, 40, 10, 100)
        layout.setSpacing(10)

        self.chat_text_edit = QTextBrowser()
        self.chat_text_edit.setReadOnly(True)
        self.chat_text_edit.setTextInteractionFlags(Qt.TextSelectableByMouse | Qt.TextSelectableByKeyboard | Qt.LinksAccessibleByMouse)
        self.chat_text_edit.setOpenExternalLinks(True)
        self.chat_text_edit.setFrameStyle(QFrame.NoFrame)
        
        # Remove any limits on content - ensure full messages are displayed
        self.chat_text_edit.document().setMaximumBlockCount(0)  # 0 = unlimited
        # Ensure document expands to fit all content
        self.chat_text_edit.document().setDocumentMargin(10)
        
        # Ensure the document can grow indefinitely to show all content
        self.chat_text_edit.setLineWrapMode(QTextEdit.WidgetWidth)
        self.chat_text_edit.setWordWrapMode(QTextOption.WrapAtWordBoundaryOrAnywhere)
        
        # Remove any size constraints that might limit content display
        size_policy = QSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        size_policy.setVerticalStretch(1)
        self.chat_text_edit.setSizePolicy(size_policy)
        
        # Set document stylesheet for futuristic link styling
        self.chat_text_edit.document().setDefaultStyleSheet("""
            a {
                color: #00dcff;
                text-decoration: none;
                border-bottom: 1px solid rgba(0, 220, 255, 0.5);
                padding: 2px 4px;
                border-radius: 4px;
                background: rgba(0, 220, 255, 0.1);
                font-weight: 500;
            }
            a:hover {
                color: #ffffff;
                background: rgba(0, 220, 255, 0.25);
                border-bottom: 1px solid rgba(0, 220, 255, 0.8);
                box-shadow: 0 0 8px rgba(0, 220, 255, 0.4);
            }
        """)
        self.chat_text_edit.setStyleSheet("""
            QTextBrowser {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1, 
                    stop:0 rgba(5, 5, 15, 0.95), 
                    stop:1 rgba(0, 0, 10, 0.98));
                color: #e0e0e0;
                border: 2px solid rgba(0, 200, 255, 0.15);
                border-radius: 20px;
                padding: 15px;
                font-size: 14px;
                font-weight: 400;
                line-height: 1.6;
                selection-background-color: rgba(0, 200, 255, 0.4);
                min-width: 300px;
                min-height: 200px;
            }
            QTextBrowser:focus {
                outline: none;
            }
            QTextBrowser:hover {
                border: 2px solid rgba(0, 200, 255, 0.25);
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1, 
                    stop:0 rgba(8, 8, 18, 0.98), 
                    stop:1 rgba(0, 0, 12, 1));
            }
            QScrollBar:vertical {
                background-color: rgba(10, 10, 20, 0.6);
                width: 14px;
                border-radius: 7px;
                border: 1px solid rgba(0, 200, 255, 0.1);
            }
            QScrollBar::handle:vertical {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 rgba(0, 200, 255, 0.6),
                    stop:1 rgba(100, 150, 255, 0.8));
                border-radius: 7px;
                min-height: 30px;
                border: 1px solid rgba(0, 200, 255, 0.3);
            }
            QScrollBar::handle:vertical:hover {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 rgba(0, 220, 255, 0.8),
                    stop:1 rgba(120, 170, 255, 1));
                box-shadow: 0 0 10px rgba(0, 200, 255, 0.5);
            }
            QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
                height: 0px;
            }
        """)
        # Add stretch factor to allow chat area to expand and show full content
        layout.addWidget(self.chat_text_edit, stretch=1)
        
        # Add GIF label after chat text edit (original position)
        self.gif_label = QLabel()
        self.gif_label.setStyleSheet("border: none; background: transparent;")
        self.gif_label.setAlignment(Qt.AlignCenter)
        gif_path = rf"{GraphicsDirPath}\FRIDAY.gif"
        if os.path.exists(gif_path):
            try:
                movie = QMovie(gif_path)
                if movie.isValid():
                    max_gif_size_W = 400
                    max_gif_size_H = 225
                    movie.setScaledSize(QSize(max_gif_size_W, max_gif_size_H))
                    movie.setCacheMode(QMovie.CacheAll)
                    self.gif_label.setMovie(movie)
                    self.gif_label.setFixedSize(max_gif_size_W, max_gif_size_H)
                    self.gif_label.setVisible(True)
                    movie.start()
                    print(f"GIF loaded successfully from: {gif_path}")
                else:
                    print(f"Warning: GIF file is not valid: {gif_path}")
                    self.gif_label.setFixedHeight(225)
            except Exception as e:
                print(f"Error loading GIF: {e}")
                self.gif_label.setText("")
                self.gif_label.setFixedHeight(225)
        else:
            self.gif_label.setText("")
            self.gif_label.setFixedHeight(225)
            print(f"Warning: GIF file not found at {gif_path}")
        layout.addWidget(self.gif_label, alignment=Qt.AlignCenter)
        
        # Add text input section
        input_layout = QHBoxLayout()
        self.text_input = QLineEdit()
        self.text_input.setPlaceholderText("Type your message here...")
        self.text_input.setStyleSheet("""
            QLineEdit {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(20, 20, 35, 0.95),
                    stop:1 rgba(15, 15, 30, 0.95));
                color: #e0e0e0;
                border: 2px solid rgba(0, 200, 255, 0.3);
                border-radius: 20px;
                padding: 12px 16px;
                font-size: 14px;
                font-weight: 500;
                min-width: 200px;
            }
            QLineEdit:hover {
                border: 2px solid rgba(0, 220, 255, 0.5);
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(25, 25, 40, 0.98),
                    stop:1 rgba(18, 18, 33, 0.98));
            }
            QLineEdit:focus {
                border: 2px solid rgba(0, 220, 255, 0.8);
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(30, 30, 45, 1),
                    stop:1 rgba(20, 20, 35, 1));
            }
        """)
        self.text_input.returnPressed.connect(self.send_text_message)
        
        # Add mode selection dropdown - Only General Assistant and Sales Assistant
        self.mode_dropdown = QComboBox()
        self.mode_dropdown.addItems([
            "General Assistant",
            "Sales Assistant"
        ])
        self.mode_dropdown.setCurrentText("General Assistant")
        self.mode_dropdown.setStyleSheet("""
            QComboBox {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(15, 25, 45, 0.95),
                    stop:1 rgba(10, 15, 35, 0.95));
                color: #e0e0e0;
                border: 2px solid rgba(0, 200, 255, 0.3);
                border-radius: 18px;
                padding: 12px 18px;
                font-size: 13px;
                font-weight: 600;
                min-width: 160px;
            }
            QComboBox:hover {
                border: 2px solid rgba(0, 220, 255, 0.6);
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(20, 30, 50, 1),
                    stop:1 rgba(15, 20, 40, 1));
            }
            QComboBox:focus {
                border: 2px solid rgba(0, 220, 255, 0.8);
            }
            QComboBox::drop-down {
                border: none;
                width: 30px;
                background: transparent;
            }
            QComboBox::down-arrow {
                image: none;
                width: 0;
                height: 0;
                border-left: 6px solid transparent;
                border-right: 6px solid transparent;
                border-top: 8px solid rgba(0, 220, 255, 0.8);
                margin-right: 8px;
            }
            QComboBox QAbstractItemView {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(10, 15, 30, 0.98),
                    stop:1 rgba(5, 10, 25, 0.98));
                color: #e0e0e0;
                border: 2px solid rgba(0, 200, 255, 0.4);
                border-radius: 12px;
                selection-background-color: rgba(0, 200, 255, 0.4);
                padding: 8px;
                font-size: 13px;
            }
            QComboBox QAbstractItemView::item {
                padding: 10px 14px;
                border-radius: 8px;
                margin: 2px;
            }
            QComboBox QAbstractItemView::item:hover {
                background-color: rgba(0, 200, 255, 0.2);
            }
        """)
        self.mode_dropdown.currentTextChanged.connect(self.on_mode_changed)
        
        # Add file upload button
        self.upload_button = QPushButton("📁")
        self.upload_button.setToolTip("Upload File")
        self.upload_button.setStyleSheet("""
            QPushButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(0, 255, 150, 0.8),
                    stop:1 rgba(0, 200, 120, 0.9));
                color: white;
                border: 2px solid rgba(0, 255, 150, 0.5);
                border-radius: 18px;
                padding: 12px 18px;
                font-size: 20px;
                font-weight: bold;
                min-width: 55px;
            }
            QPushButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(0, 255, 180, 0.95),
                    stop:1 rgba(0, 220, 140, 1));
                border: 2px solid rgba(0, 255, 180, 0.8);
            }
            QPushButton:pressed {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(0, 220, 130, 0.9),
                    stop:1 rgba(0, 180, 110, 1));
            }
        """)
        self.upload_button.clicked.connect(self.upload_file)
        
        self.send_button = QPushButton("Send")
        self.send_button.setStyleSheet("""
            QPushButton {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 rgba(0, 150, 255, 0.9),
                    stop:1 rgba(100, 50, 255, 0.9));
                color: white;
                border: 2px solid rgba(0, 200, 255, 0.5);
                border-radius: 18px;
                padding: 14px 28px;
                font-size: 15px;
                font-weight: 700;
                letter-spacing: 0.5px;
            }
            QPushButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 rgba(0, 180, 255, 1),
                    stop:1 rgba(120, 70, 255, 1));
                border: 2px solid rgba(0, 220, 255, 0.8);
            }
            QPushButton:pressed {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 rgba(0, 130, 230, 0.95),
                    stop:1 rgba(80, 40, 230, 0.95));
            }
        """)
        self.send_button.clicked.connect(self.send_text_message)
        
        # Add Reset button
        self.reset_button = QPushButton("Reset")
        self.reset_button.setToolTip("Reset assistant to Available state")
        self.reset_button.setStyleSheet("""
            QPushButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(255, 150, 0, 0.8),
                    stop:1 rgba(255, 100, 0, 0.9));
                color: white;
                border: 2px solid rgba(255, 180, 0, 0.5);
                border-radius: 18px;
                padding: 12px 20px;
                font-size: 14px;
                font-weight: 700;
            }
            QPushButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(255, 180, 0, 0.95),
                    stop:1 rgba(255, 120, 0, 1));
                border: 2px solid rgba(255, 200, 0, 0.8);
            }
            QPushButton:pressed {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 rgba(255, 130, 0, 0.9),
                    stop:1 rgba(230, 90, 0, 1));
            }
        """)
        self.reset_button.clicked.connect(self.reset_assistant)
        
        # Responsive input layout
        input_layout.addWidget(self.text_input, 3)  # Give text input more space
        input_layout.addWidget(self.mode_dropdown, 0)
        input_layout.addWidget(self.upload_button, 0)
        input_layout.addWidget(self.reset_button, 0)
        input_layout.addWidget(self.send_button, 0)
        input_layout.setSpacing(8)  # Add spacing between input elements
        layout.addLayout(input_layout)

        self.setStyleSheet("""
            QWidget {
                background-color: #000000;
            }
        """)
        layout.setSizeConstraint(QVBoxLayout.SetDefaultConstraint)
        layout.setStretch(1, 1)
        self.setSizePolicy(QSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding))

        text_color = QColor(Qt.blue)
        text_color_text = QTextCharFormat()
        text_color_text.setForeground(text_color)
        self.chat_text_edit.setCurrentCharFormat(text_color_text)

        self.label = QLabel("")
        self.label.setStyleSheet("""
            color: #e0e0e0; 
            font-size: 15px; 
            font-weight: 600;
            border: none; 
            margin-top: -30px;
            padding: 10px 20px;
            background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                stop:0 rgba(0, 200, 255, 0.15),
                stop:1 rgba(100, 150, 255, 0.15));
            border-radius: 25px;
            border: 2px solid rgba(0, 200, 255, 0.4);
            letter-spacing: 0.5px;
        """)
        self.label.setAlignment(Qt.AlignRight)
        layout.addWidget(self.label)

        font = QFont()
        font.setPointSize(13)
        self.chat_text_edit.setFont(font)

        self.timer = QTimer(self)
        self.timer.timeout.connect(self.loadMessages)
        self.timer.timeout.connect(self.SpeechRecogText)
        self.timer.start(10)  # Ultra fast - 10ms for maximum responsiveness
        
        # Initialize displayed messages tracking to prevent duplicates
        self._displayed_messages = set()
        
        # Add pulsing animation timer for status label
        self.pulse_timer = QTimer(self)
        self.pulse_timer.timeout.connect(self.pulse_status_label)
        self.pulse_timer.start(2000)  # Pulse every 2 seconds
        self.pulse_state = 0

        self.chat_text_edit.viewport().installEventFilter(self)
        self.setStyleSheet("""
            QScrollBar:vertical {
                border: none;
                background: black;
                width: 10px;
                margin: 0px 0px 0px 0px;
            }

            QScrollBar::handle:vertical {
                background: white;
                min-height: 20px;
            }

            QScrollBar::add-line:vertical {
                background: black;
                subcontrol-position: bottom;
                subcontrol-origin: margin;
                height: 10px;
            }

            QScrollBar::sub-line:vertical {
                background: black;
                subcontrol-position: top;
                subcontrol-origin: margin;
                height: 10px;
            }

            QScrollBar::up-arrow:vertical, QScrollBar::down-arrow:vertical {
                border: none;
                background: none;
                color: none;
            }

            QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical {
                background: none;
            }

        """)

    def loadMessages(self):
        global old_chat_message
        try:
            with open(rf'{TempDirPath}\Responses.data', 'r', encoding='utf-8') as file:
                messages = file.read()
            
            # Only process if content has changed
            if messages and messages.strip() and messages != old_chat_message:
                assistant_name = env_vars.get('Assistantname', 'JARVIS')
                
                # Parse messages without clearing to prevent duplication
                # Then parse messages line by line, ensuring each sender gets a separate bubble
                lines = messages.split('\n')
                parsed_messages = []
                current_message = ""
                current_sender = None
                
                for line in lines:
                    line_original = line  # Keep original for preserving newlines
                    line = line.strip()
                    
                    # Check for sender patterns at the START of line
                    line_lower = line.lower()
                    
                    # User patterns (case-insensitive detection)
                    if line_lower.startswith("shubh mishra:") or (":" in line and any(name in line_lower[:30] for name in ["shubh mishra:", "you:", "user:"])):
                        # Finalize previous message if exists
                        if current_message and current_sender:
                            parsed_messages.append((current_sender, current_message.strip()))
                            current_message = ""
                        # Determine sender name
                        if "shubh" in line_lower or "mishra" in line_lower:
                            current_sender = "Shubh Mishra"
                        else:
                            current_sender = "You"
                        # Extract message content after colon, preserve leading/trailing whitespace from original
                        if ":" in line:
                            current_message = line.split(":", 1)[1]
                            # If original line had more content after the colon, preserve it
                            if ":" in line_original:
                                colon_pos = line_original.find(":")
                                current_message = line_original[colon_pos+1:]
                        else:
                            current_message = ""
                    # Assistant patterns (case-insensitive)
                    elif line_lower.startswith("jarvis:") or line_lower.startswith(f"{assistant_name.lower()}:"):
                        # Finalize previous message if exists
                        if current_message and current_sender:
                            parsed_messages.append((current_sender, current_message.strip()))
                            current_message = ""
                        if line_original.startswith(f"{assistant_name}:"):
                            current_sender = assistant_name
                            if ":" in line_original:
                                colon_pos = line_original.find(":")
                                current_message = line_original[colon_pos+1:]
                            else:
                                current_message = ""
                        else:
                            current_sender = "JARVIS"
                            if ":" in line:
                                if ":" in line_original:
                                    colon_pos = line_original.find(":")
                                    current_message = line_original[colon_pos+1:]
                                else:
                                    current_message = line.split(":", 1)[1]
                            else:
                                current_message = ""
                    else:
                        # Continuation of current message - preserve newlines and full content
                        if current_sender:
                            if current_message:
                                # Always append with newline to preserve formatting
                                current_message += "\n" + line_original
                            else:
                                current_message = line_original
                        # If no sender yet, this might be a standalone line or formatting - skip it
                
                # Add final message if exists
                if current_message and current_sender:
                    parsed_messages.append((current_sender, current_message.strip()))
                
                # Track displayed messages to prevent duplicates
                displayed_count = 0
                for sender, message_content in parsed_messages:
                    if message_content:  # Only add non-empty messages
                        # Check if this exact message was already displayed (prevent duplicates)
                        message_key = f"{sender}:{message_content[:50]}"  # Use first 50 chars as key
                        if not hasattr(self, '_displayed_messages'):
                            self._displayed_messages = set()
                        
                        if message_key not in self._displayed_messages:
                            print(f"Loading message from {sender}: {len(message_content)} chars, preview: {message_content[:100]}...")
                            self.addMessage(f"{sender}: {message_content}", color='White')
                            self._displayed_messages.add(message_key)
                            displayed_count += 1
                            # Keep only last 20 displayed messages to prevent memory growth
                            if len(self._displayed_messages) > 20:
                                # Remove oldest entries (simple: clear and re-add recent ones)
                                self._displayed_messages = set(list(self._displayed_messages)[-10:])
                
                # Only update old_chat_message and clear file if we actually displayed new messages
                if displayed_count > 0:
                    old_chat_message = messages
                    # Clear the file after displaying to prevent re-display
                    try:
                        with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as file:
                            file.write('')  # Clear the file
                    except Exception as e:
                        print(f"Error clearing Responses.data: {e}")
                
                # Force scroll to bottom after loading messages - delayed to ensure content is rendered
                QTimer.singleShot(100, lambda: self._scroll_to_bottom(10))
                QTimer.singleShot(300, lambda: self._scroll_to_bottom(11))
                QTimer.singleShot(600, lambda: self._scroll_to_bottom(12))
        except FileNotFoundError:
            pass
        except Exception as e:
            print(f"Error parsing messages: {e}")
            # Fallback: try to add as single message
            try:
                if messages and messages != old_chat_message:
                    self.addMessage(message=messages, color='White')
                    old_chat_message = messages
            except:
                pass

    def SpeechRecogText(self):
        try:
            with open(rf'{TempDirPath}\Status.data', 'r', encoding='utf-8') as file:
                messages = file.read()
            self.label.setText(messages)
        except FileNotFoundError:
            pass

    def load_icon(self, path, width=60, height=60):
        pixmap = QPixmap(path)
        new_pixmap = pixmap.scaled(width, height)
        self.icon_label.setPixmap(new_pixmap)

    def toggle_icon(self, event=None):
        if self.toggled:
            # Mic is currently ON, turning it OFF
            self.load_icon(rf'{GraphicsDirPath}\Mic_off.png', 60, 60)
            MicButtonClosed()
            self.toggled = False
        else:
            # Mic is currently OFF, turning it ON
            self.load_icon(rf'{GraphicsDirPath}\Mic_on.png', 60, 60)
            MicButtonInitiated()
            self.toggled = True

    def addMessage(self, message, color):
        # Clean up any raw HTML/CSS code that might be in the message
        import re
        import html
        
        # Remove any raw HTML style attributes that appear as text (like '." style="..."')
        # Pattern to match things like '." style="...">' that shouldn't be in the text
        processed_message = message
        
        # First, clean up any broken HTML that might be in the text
        # Remove standalone style attributes that aren't part of proper HTML tags
        processed_message = re.sub(r'\.?\s*style="[^"]*"[^>]*>', '', processed_message)
        processed_message = re.sub(r'\.?\s*style=\'[^\']*\'[^>]*>', '', processed_message)
        
        # Check for markdown links and convert to clickable HTML links
        link_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
        
        def replace_link(match):
            link_text = match.group(1)
            link_url = match.group(2)
            # Escape the link text to prevent HTML injection
            link_text_escaped = html.escape(link_text)
            return f'<a href="{link_url}" style="color: #00dcff; text-decoration: none; border-bottom: 1px solid rgba(0, 220, 255, 0.5); padding: 2px 4px; border-radius: 4px; background: rgba(0, 220, 255, 0.1); transition: all 0.3s ease; font-weight: 500;">{link_text_escaped}</a>'
        
        # Replace markdown links with HTML links
        processed_message = re.sub(link_pattern, replace_link, processed_message)
        
        # Also detect and style plain URLs (http://, https://, www.)
        # Make sure we don't match URLs that are already inside HTML tags
        url_pattern = r'(?<!href=["\'])(?<!["\']>)(https?://[^\s<>"\']+|www\.[^\s<>"\']+)(?![^<]*</a>)'
        def replace_url(match):
            url = match.group(0).rstrip('.,;!?')
            display_url = url
            # Make URL clickable
            return f'<a href="{url if url.startswith("http") else "https://" + url}" style="color: #00dcff; text-decoration: none; border-bottom: 1px solid rgba(0, 220, 255, 0.5); padding: 2px 4px; border-radius: 4px; background: rgba(0, 220, 255, 0.1); font-weight: 500;">{display_url}</a>'
        processed_message = re.sub(url_pattern, replace_url, processed_message)
        
        # Remove any remaining broken HTML fragments that might have been left
        processed_message = re.sub(r'\.?\s*style=["\'][^"\']*["\']\s*>', '', processed_message)
        
        # Clean up any orphaned HTML fragments (like style attributes without tags)
        # Remove patterns like '." style="..."' or just 'style="..."'
        processed_message = re.sub(r'\.\s*["\']?\s*style\s*=\s*["\'][^"\']*["\']\s*[>]?', '.', processed_message)
        processed_message = re.sub(r'\s*style\s*=\s*["\'][^"\']*["\']\s*[>]?', '', processed_message)
        
        # Extract our properly formatted link tags before escaping
        link_pattern = r'<a\s+[^>]*href=["\'][^"\']*["\'][^>]*>.*?</a>'
        link_matches = list(re.finditer(link_pattern, processed_message, re.IGNORECASE | re.DOTALL))
        link_placeholders = []
        
        # Replace links with placeholders
        offset = 0
        for i, match in enumerate(link_matches):
            placeholder = f"__LINK_PLACEHOLDER_{i}__"
            link_placeholders.append(match.group(0))
            start, end = match.span()
            processed_message = processed_message[:start+offset] + placeholder + processed_message[end+offset:]
            offset += len(placeholder) - (end - start)
        
        # Escape HTML in the remaining text (but not our links)
        processed_message = html.escape(processed_message)
        
        # Restore our properly formatted link tags
        for i, link_tag in enumerate(link_placeholders):
            processed_message = processed_message.replace(f"__LINK_PLACEHOLDER_{i}__", link_tag)
        
        # Final cleanup - remove any remaining broken fragments
        processed_message = re.sub(r'\.\s*["\']?\s*style\s*=\s*["\'][^"\']*["\']\s*[>]?', '.', processed_message)
        
        # Determine if this is a user message or assistant message
        is_user_message = False
        is_assistant_message = False
        
        # Get assistant name from environment
        assistant_name = env_vars.get('Assistantname', 'JARVIS')
        
        # Check message content to determine sender (enhanced detection)
        sender = None
        message_content = processed_message
        
        # Extract sender and content from message
        if ":" in processed_message:
            parts = processed_message.split(":", 1)
            if len(parts) > 1:
                potential_sender = parts[0].strip()
                message_content = parts[1].strip()
                
                # Check if it's a user message (case-insensitive check)
                potential_sender_lower = potential_sender.lower()
                if (potential_sender_lower in ["shubh mishra", "you", "user", "your name"] or 
                    "shubh" in potential_sender_lower or "mishra" in potential_sender_lower):
                    is_user_message = True
                    if "shubh" in potential_sender_lower or "mishra" in potential_sender_lower:
                        sender = "Shubh Mishra"
                    else:
                        sender = "You"
                # Check if it's an assistant message (case-insensitive)
                elif potential_sender_lower in ["jarvis", assistant_name.lower(), "assistant"] or potential_sender == assistant_name:
                    is_assistant_message = True
                    if potential_sender == assistant_name:
                        sender = assistant_name
                    else:
                        sender = "JARVIS"
        
        # Fallback detection if colon parsing didn't work (case-insensitive)
        if not sender:
            message_lower = message.lower()
            if message_lower.startswith("shubh mishra:") or message_lower.startswith("you:") or "shubh" in message_lower[:15] or "your name:" in message_lower:
                is_user_message = True
                sender = "Shubh Mishra" if "shubh" in message_lower or "mishra" in message_lower else "You"
            elif message_lower.startswith("jarvis:") or message_lower.startswith(f"{assistant_name.lower()}:") or message.startswith(f"{assistant_name}:"):
                is_assistant_message = True
                sender = assistant_name if message.startswith(f"{assistant_name}:") else "JARVIS"
        
        # Create styled message bubble based on sender - ALL messages LEFT aligned
        if is_user_message:
            # User message - LEFT aligned (changed from right), futuristic gradient blue
            html_message = f"""
            <div style="margin: 0px 0px 25px 0px; padding: 0px; clear: both; display: block; width: 100%; box-sizing: border-box;">
                <div style="margin-left: 0px; margin-right: 10px; max-width: 85%; min-width: 200px; text-align: left; display: block;">
                    <div style="display: inline-block; 
                        background: linear-gradient(135deg, rgba(0, 150, 255, 0.9) 0%, rgba(100, 50, 255, 0.9) 100%);
                        color: #ffffff; 
                        padding: 14px 18px; 
                        border-radius: 20px 20px 20px 6px; 
                        word-wrap: break-word; 
                        word-break: break-word;
                        overflow-wrap: break-word;
                        box-shadow: 0 4px 15px rgba(0, 150, 255, 0.4), 
                                    0 0 20px rgba(100, 50, 255, 0.2),
                                    inset 0 1px 0 rgba(255, 255, 255, 0.2);
                        border: 1px solid rgba(255, 255, 255, 0.2);
                        border-left: 4px solid rgba(0, 150, 255, 0.8);
                        backdrop-filter: blur(10px);
                        position: relative;
                        width: fit-content;
                        max-width: 100%;
                        box-sizing: border-box;">
                        <strong style="font-size: 12px; opacity: 1; font-weight: 700; letter-spacing: 1px; text-transform: uppercase; display: block; margin-bottom: 8px; color: #ffffff; text-shadow: 0 0 8px rgba(255, 255, 255, 0.6);">{sender if sender else "You"}</strong>
                        <div style="line-height: 1.6; font-size: 14.5px; display: block; word-wrap: break-word; word-break: break-word; white-space: pre-wrap; overflow-wrap: break-word; hyphens: auto; max-height: none; overflow: visible;">{message_content}</div>
                    </div>
                </div>
            </div>
            """
        elif is_assistant_message:
            # Assistant message - left aligned, futuristic dark with glow, list-like
            html_message = f"""
            <div style="margin: 0px 0px 25px 0px; padding: 0px; clear: both; display: block; width: 100%; box-sizing: border-box;">
                <div style="margin-left: 0px; margin-right: 10px; max-width: 85%; min-width: 200px; text-align: left; display: block;">
                    <div style="display: inline-block;
                        background: linear-gradient(135deg, rgba(20, 25, 40, 0.95) 0%, rgba(10, 15, 30, 0.98) 100%);
                        color: #e0e0e0; 
                        padding: 14px 18px; 
                        border-radius: 20px 20px 20px 6px; 
                        word-wrap: break-word; 
                        word-break: break-word;
                        overflow-wrap: break-word;
                        box-shadow: 0 4px 15px rgba(0, 200, 255, 0.3), 
                                    0 0 25px rgba(0, 200, 255, 0.15),
                                    inset 0 1px 0 rgba(255, 255, 255, 0.1);
                        border: 1px solid rgba(0, 200, 255, 0.4);
                        border-left: 4px solid rgba(0, 220, 255, 0.8);
                        backdrop-filter: blur(10px);
                        position: relative;
                        width: fit-content;
                        max-width: 100%;
                        box-sizing: border-box;">
                        <strong style="font-size: 12px; color: #00dcff; font-weight: 700; letter-spacing: 1px; text-transform: uppercase; text-shadow: 0 0 12px rgba(0, 220, 255, 0.8), 0 0 20px rgba(0, 220, 255, 0.4); display: block; margin-bottom: 8px; border-bottom: 1px solid rgba(0, 220, 255, 0.3); padding-bottom: 4px;">{sender if sender else assistant_name}</strong>
                        <div style="line-height: 1.7; font-size: 14.5px; display: block; word-wrap: break-word; word-break: break-word; white-space: pre-wrap; overflow-wrap: break-word; hyphens: auto; max-height: none; overflow: visible;">{message_content}</div>
                    </div>
                </div>
            </div>
            """
        else:
            # Generic message - use color parameter with futuristic styling
            color_map = {
                "White": "#e0e0e0",
                "LightGreen": "#00ff88",
                "LightBlue": "#00dcff",
                "Green": "#00ff00",
                "Red": "#ff4444",
                "Yellow": "#ffd700",
                "LightYellow": "#ffff99"
            }
            text_color = color_map.get(color, "#e0e0e0")
            # Convert hex to rgba for glow effect
            hex_color = text_color.lstrip('#')
            rgb_tuple = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            rgba_glow = f"rgba({rgb_tuple[0]}, {rgb_tuple[1]}, {rgb_tuple[2]}, 0.15)"
            html_message = f"""
            <div style="margin: 0px 0px 25px 0px; padding: 0px; clear: both; display: block;">
                <div style="margin-left: 0px; max-width: 70%; text-align: left; display: block;">
                    <div style="display: inline-block;
                        background: linear-gradient(135deg, rgba(20, 25, 40, 0.95) 0%, rgba(10, 15, 30, 0.98) 100%);
                        color: {text_color}; 
                        padding: 14px 18px; 
                        border-radius: 20px; 
                        word-wrap: break-word; 
                        box-shadow: 0 4px 15px rgba(0, 200, 255, 0.2), 
                                    0 0 20px {rgba_glow},
                                    inset 0 1px 0 rgba(255, 255, 255, 0.1);
                        border: 1px solid rgba(0, 200, 255, 0.3);
                        border-left: 4px solid {text_color};
                        backdrop-filter: blur(10px);
                        width: fit-content;
                        max-width: 100%;">
                        <div style="line-height: 1.7; font-size: 14.5px; display: block; word-wrap: break-word; word-break: break-word; white-space: pre-wrap; overflow-wrap: break-word; hyphens: auto; max-height: none; overflow: visible;">{message_content}</div>
                    </div>
                </div>
            </div>
            """
        
        # Insert HTML message with proper line breaks - ensure full content is displayed
        # Use insertHtml to append to existing content (prevents duplication)
        cursor = self.chat_text_edit.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        # Add clear line break before new message for list-like appearance
        if self.chat_text_edit.toPlainText().strip():
            cursor.insertHtml("<br>")
        
        # Insert the full HTML message - ensure it's complete
        cursor.insertHtml(html_message)
        cursor.insertBlock()  # Ensure new message starts on new line
        
        # Force update the document to render all HTML content
        self.chat_text_edit.document().setModified(True)
        
        # Force the document to update its size to accommodate all content
        try:
            # This ensures the document recalculates its size
            doc = self.chat_text_edit.document()
            # Force layout recalculation
            doc.adjustSize()
            # Update the viewport
            self.chat_text_edit.viewport().update()
            # Force repaint
            self.chat_text_edit.repaint()
            print(f"Document size after insert: {doc.size().width()}x{doc.size().height()}")
        except Exception as e:
            print(f"Error adjusting document size: {e}")
        
        # Move cursor to end
        cursor = self.chat_text_edit.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_text_edit.setTextCursor(cursor)
        
        # Process events immediately to let UI update
        app = QApplication.instance()
        if app:
            app.processEvents()
        
        # Force scroll to bottom to show the latest content - multiple attempts to ensure it works
        scrollbar = self.chat_text_edit.verticalScrollBar()
        
        # Method 1: Scroll to maximum
        max_scroll = scrollbar.maximum()
        scrollbar.setValue(max_scroll)
        
        # Process events to let scrollbar update
        if app:
            app.processEvents()
        
        # Method 2: Ensure cursor is visible (forces scroll)
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_text_edit.setTextCursor(cursor)
        self.chat_text_edit.ensureCursorVisible()
        
        # Process events again
        if app:
            app.processEvents()
        
        # Method 3: Scroll to maximum again after events processed
        max_scroll = scrollbar.maximum()
        scrollbar.setValue(max_scroll)
        
        # Method 4: Use movePosition to end and ensure visible again
        cursor = self.chat_text_edit.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_text_edit.setTextCursor(cursor)
        self.chat_text_edit.ensureCursorVisible()
        
        # Final scroll to maximum
        scrollbar.setValue(scrollbar.maximum())
        
        # Force a repaint to ensure all content is displayed
        self.chat_text_edit.repaint()
        self.chat_text_edit.update()
        
        # Process events one more time to ensure UI updates
        if app:
            app.processEvents()
        
        # Use QTimer with multiple attempts to ensure scrolling happens after content is fully rendered
        # These delayed scrolls catch content that loads asynchronously
        QTimer.singleShot(50, lambda: self._scroll_to_bottom(1))
        QTimer.singleShot(150, lambda: self._scroll_to_bottom(2))
        QTimer.singleShot(300, lambda: self._scroll_to_bottom(3))
        QTimer.singleShot(500, lambda: self._scroll_to_bottom(4))
        QTimer.singleShot(1000, lambda: self._scroll_to_bottom(5))  # Extra attempt for very long messages
        QTimer.singleShot(1500, lambda: self._scroll_to_bottom(6))  # Even longer delay for very long content
        QTimer.singleShot(2000, lambda: self._scroll_to_bottom(7))  # Maximum delay for complete rendering
    
    def _scroll_to_bottom(self, attempt=1):
        """Helper method to scroll to bottom after content is rendered - multiple attempts"""
        try:
            scrollbar = self.chat_text_edit.verticalScrollBar()
            
            # Get document and viewport heights
            doc = self.chat_text_edit.document()
            viewport_height = self.chat_text_edit.viewport().height()
            doc_height = doc.size().height()
            
            # Calculate the maximum scroll value based on actual content
            max_scroll = max(0, int(doc_height - viewport_height))
            
            # Set scrollbar range if needed
            if scrollbar.maximum() < max_scroll:
                scrollbar.setMaximum(max_scroll)
            
            scrollbar.setValue(max_scroll)
            
            # Also ensure cursor is at end
            cursor = self.chat_text_edit.textCursor()
            cursor.movePosition(cursor.MoveOperation.End)
            self.chat_text_edit.setTextCursor(cursor)
            self.chat_text_edit.ensureCursorVisible()
            
            # Process events to ensure scroll happens
            app = QApplication.instance()
            if app:
                app.processEvents()
            
            # Force another scroll after events processed
            new_max = scrollbar.maximum()
            if new_max > 0:
                scrollbar.setValue(new_max)
            
            # Additional scroll attempts with different methods
            # Sometimes the scrollbar maximum needs to be recalculated
            doc_size = doc.size()
            if doc_size.height() > viewport_height:
                # Force update the scrollbar range
                scrollbar.setRange(0, max(0, int(doc_size.height() - viewport_height)))
                scrollbar.setValue(scrollbar.maximum())
            
            print(f"Scroll attempt {attempt}: doc_height={doc_height:.1f}, viewport={viewport_height}, max_scroll={max_scroll}, scrollbar_max={scrollbar.maximum()}, scrollbar_value={scrollbar.value()}")
        except Exception as e:
            import traceback
            print(f"Scroll attempt {attempt} error: {e}")
            print(traceback.format_exc())
    
    def reset_assistant(self):
        """Reset the assistant and clear ALL chats"""
        try:
            print("Reset button clicked - clearing all chats and resetting")
            
            # Stop any ongoing speech
            try:
                from Backend.TextToSpeech import interrupt_speech
                interrupt_speech()
            except:
                pass
            
            # Reset all states
            try:
                from Main import safe_interrupt_speech
                safe_interrupt_speech()
            except:
                pass
            
            # Check mic status BEFORE clearing (to preserve state)
            mic_was_on = False
            try:
                from Frontend.GUI import GetMicrophoneStatus
                mic_status = GetMicrophoneStatus()
                mic_was_on = mic_status and mic_status.lower() == "true"
                print(f"Reset: Mic status before reset: {mic_status} (was_on: {mic_was_on})")
            except Exception as e:
                print(f"Error checking mic status before reset: {e}")
            
            # CLEAR ALL CHATS - Clear the chat display
            self.chat_text_edit.clear()
            
            # Clear chat data files
            try:
                with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as f:
                    f.write("")
                with open(rf'{TempDirPath}\Database.data', 'w', encoding='utf-8') as f:
                    f.write("")
                import json
                with open(r'Data\ChatLog.json', 'w', encoding='utf-8') as f:
                    json.dump([], f)
            except Exception as e:
                print(f"Error clearing chat files: {e}")
            
            # Clear conversation memory
            try:
                from Backend.Memory import clear_conversation_memory
                clear_conversation_memory()
            except Exception as e:
                print(f"Error clearing memory: {e}")
            
            # Set status based on microphone status - keep listening if mic was ON
            # Don't turn mic off during reset - preserve the user's mic state
            if mic_was_on:
                # Mic was ON - keep it ON and set to Listening
                try:
                    SetAsssistantStatus("Listening...")
                    self.label.setText("Listening...")
                    print("Reset: Mic was ON, keeping status as Listening...")
                except Exception as e:
                    print(f"Error setting Listening status: {e}")
            else:
                # Mic was OFF - set to Available
                try:
                    SetAsssistantStatus("Available...")
                    self.label.setText("Available...")
                    print("Reset: Mic was OFF, setting status to Available...")
                except Exception as e:
                    print(f"Error setting Available status: {e}")
            
            # Reset global old_chat_message
            global old_chat_message
            old_chat_message = ""
            
            # Clear displayed messages tracking to prevent duplicates after reset
            if hasattr(self, '_displayed_messages'):
                self._displayed_messages.clear()
            
            # Show reset confirmation
            self.addMessage("JARVIS: All chats cleared! I'm ready to help!", "LightGreen")
            print("Reset completed - all chats cleared")
            
        except Exception as e:
            print(f"Error in reset: {e}")
            self.addMessage("JARVIS: Reset completed with minor issues. Ready to help!", "LightGreen")
    
    def clear_chat_history(self):
        """Clear only the chat history (display and files) without full reset"""
        try:
            print("Clearing chat history...")
            
            # Clear the chat display
            self.chat_text_edit.clear()
            
            # Clear chat data files
            try:
                with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as f:
                    f.write("")
                with open(rf'{TempDirPath}\Database.data', 'w', encoding='utf-8') as f:
                    f.write("")
                import json
                with open(r'Data\ChatLog.json', 'w', encoding='utf-8') as f:
                    json.dump([], f)
            except Exception as e:
                print(f"Error clearing chat files: {e}")
            
            # Clear conversation memory
            try:
                from Backend.Memory import clear_conversation_memory
                clear_conversation_memory()
            except Exception as e:
                print(f"Error clearing memory: {e}")
            
            # Reset global old_chat_message
            global old_chat_message
            old_chat_message = ""
            
            # Clear displayed messages tracking to prevent duplicates
            if hasattr(self, '_displayed_messages'):
                self._displayed_messages.clear()
            if hasattr(self, '_processing_messages'):
                self._processing_messages.clear()
            if hasattr(self, '_last_displayed_answer'):
                self._last_displayed_answer = None
            
            # Show confirmation
            self.addMessage("JARVIS: Chat history cleared!", "LightGreen")
            print("Chat history cleared successfully")
            
        except Exception as e:
            print(f"Error clearing chat history: {e}")
            import traceback
            print(traceback.format_exc())
    
    def send_text_message(self):
        """Handle text message sending"""
        try:
            message = self.text_input.text().strip()
            print("=" * 50)
            print(f"SEND BUTTON CLICKED - Text input received: '{message}'")
            print("=" * 50)
            
            if message:
                # IMMEDIATELY interrupt any ongoing speech when new message is sent
                try:
                    from Backend.TextToSpeech import interrupt_speech
                    from Main import force_interrupt_speech
                    import pygame
                    
                    # Force interrupt immediately
                    print("New question received - immediately interrupting speech")
                    interrupt_speech()
                    force_interrupt_speech()
                    
                    # Set flags to ensure speech stops
                    import Backend.TextToSpeech as tts_module
                    tts_module.speech_interrupted = True
                    tts_module.interrupt_requested = True
                    
                    # Also stop any audio playback immediately
                    try:
                        if pygame.mixer.get_init() is not None:
                            pygame.mixer.music.stop()
                            pygame.mixer.quit()
                    except:
                        pass
                    
                except Exception as e:
                    print(f"Error interrupting speech: {e}")
                
                # Clear the input field
                self.text_input.clear()
                
                # Add user message to chat (only if not already added by voice input)
                # Check if this message was already processed via voice
                from Main import last_processed_query, last_processed_time
                import time
                import re
                message_normalized = message.strip().lower()
                message_normalized = re.sub(r'[^\w\s]', '', message_normalized)
                message_normalized = ' '.join(message_normalized.split())
                current_time = time.time()
                
                # ALWAYS add user message to chat FIRST (before processing) - ensure it's visible
                # Check if it's a duplicate from voice input
                is_duplicate = (last_processed_query == message_normalized and (current_time - last_processed_time) < 2.0)
                
                # ALWAYS display user message, even if duplicate (user needs to see their input)
                print(f"GUI: Adding user message to chat: '{message}'")
                self.addMessage(f"You: {message}", "White")
                self.chat_text_edit.ensureCursorVisible()  # Ensure message is visible
                self.update()  # Force immediate GUI update
                QApplication.processEvents()  # Process events to ensure display
                
                if not is_duplicate:
                    # Also write user message to file to maintain order
                    try:
                        current_content = ""
                        with open(rf'{TempDirPath}\Responses.data', 'r', encoding='utf-8') as f:
                            current_content = f.read()
                    except:
                        pass
                    # Append user message to file content
                    if current_content and not current_content.endswith('\n'):
                        current_content += '\n'
                    current_content += f"You: {message}\n"
                    with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as f:
                        f.write(current_content)
                else:
                    print(f"GUI: Skipping duplicate user message (already processed via voice): '{message}'")
                
                # Small delay to ensure speech interruption is processed
                import time
                time.sleep(0.1)
                
                # Process the message through the AI system
                print(f"Processing text message: {message}")  # Debug
                try:
                    self.process_text_message(message)
                    print(f"Successfully called process_text_message for: '{message}'")
                except Exception as e:
                    print(f"ERROR in process_text_message: {e}")
                    import traceback
                    print(traceback.format_exc())
                    self.addMessage("JARVIS: I encountered an error processing your message. Please try again.", "Red")
            else:
                print("Empty message received")  # Debug
        except Exception as e:
            print(f"ERROR in send_text_message: {e}")
            import traceback
            print(traceback.format_exc())
    
    def process_text_message(self, message):
        """Process text message through the AI system"""
        # Track if we've already generated a response for this message to prevent duplicates
        if not hasattr(self, '_processing_messages'):
            self._processing_messages = set()
        
        # Create a unique key for this message
        import time
        message_key = f"{message.lower().strip()}_{int(time.time())}"
        
        # Check if we're already processing this exact message
        if message_key in self._processing_messages:
            print(f"Already processing this message, skipping duplicate: '{message}'")
            return
        
        # Mark as processing
        self._processing_messages.add(message_key)
        # Clean up old entries (keep only last 10)
        if len(self._processing_messages) > 10:
            self._processing_messages = set(list(self._processing_messages)[-10:])
        
        try:
            # Force interrupt any current speech when new text input is received
            from Main import force_interrupt_speech
            force_interrupt_speech()
            
            # Clear memory if user asks to start fresh
            if "clear memory" in message.lower() or "forget everything" in message.lower():
                from Backend.Memory import clear_conversation_memory
                clear_conversation_memory()
                self.addMessage("JARVIS: Memory cleared! Starting fresh.", "LightBlue")
                self.label.setText("Available...")
                return
            
            # Clear chat history if user requests it
            if any(phrase in message.lower() for phrase in ["clear chat", "clear chat history", "clear history", "clear all chats", "delete chat history"]):
                self.clear_chat_history()
                return
            
            # Check for Google Drive links (handle both plain links and HTML-wrapped links)
            import re
            # First, clean HTML tags from message
            message_clean = re.sub(r'<[^>]+>', '', message)
            # Pattern to match Drive links (more flexible)
            drive_link_pattern = r'https?://(?:drive\.google\.com|docs\.google\.com)[^\s<>"]+'
            drive_links = re.findall(drive_link_pattern, message_clean)
            
            # Also try to extract from HTML href attributes
            if not drive_links:
                href_pattern = r'href=["\'](https?://(?:drive\.google\.com|docs\.google\.com)[^\s<>"\']+)["\']'
                drive_links = re.findall(href_pattern, message)
            
            if drive_links:
                # Process Drive links
                for drive_link in drive_links:
                    # Clean the link (remove any trailing characters)
                    drive_link = drive_link.strip().rstrip('>').rstrip('"').rstrip("'")
                    
                    self.addMessage(f"🔗 Processing Google Drive link...", "LightBlue")
                    self.label.setText("Processing Drive link...")
                    self.update()
                    
                    try:
                        from Backend.DriveProcessor import process_drive_link
                        print(f"Processing Drive link: {drive_link}")
                        result = process_drive_link(drive_link)
                        print(f"Processing result: {result}")
                        print(f"Result type: {type(result)}")
                        if isinstance(result, dict):
                            print(f"Result keys: {list(result.keys())}")
                            print(f"Success value: {result.get('success')}")
                            print(f"Error value: {result.get('error')}")
                        
                        if result and isinstance(result, dict) and result.get("success"):
                            files_count = result.get("files_processed", 0)
                            entries_count = result.get("entries_created", 0)
                            processed_files = result.get("processed_files", [])
                            
                            success_msg = f"✅ Google Drive link processed successfully!\n"
                            success_msg += f"📁 Files processed: {files_count}\n"
                            success_msg += f"📝 Knowledge entries created: {entries_count}\n"
                            if processed_files:
                                # Show actual filenames if available, otherwise show count
                                file_display = [f[:30] + "..." if len(f) > 30 else f for f in processed_files[:5]]
                                success_msg += f"📄 Sample files: {', '.join(file_display)}"
                                if len(processed_files) > 5:
                                    success_msg += f" and {len(processed_files) - 5} more..."
                            
                            if result.get('note'):
                                success_msg += f"\nℹ️ {result.get('note')}"
                            
                            self.addMessage(success_msg, "LightGreen")
                            
                            # Clear, simple success message
                            success_announcement = f"Google Drive link processed successfully! {files_count} files processed with {entries_count} knowledge entries created. You can now ask me any questions about the files."
                            self.addMessage(f"JARVIS: {success_announcement}", "White")
                            self.label.setText("Available...")
                            
                            # Announce verbally in a separate thread
                            def announce_success_verbally():
                                try:
                                    import time
                                    time.sleep(0.5)  # Small delay to let GUI update
                                    from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt, TextToSpeech
                                    reset_speech_interrupt()
                                    
                                    # Short verbal announcement
                                    verbal_msg = f"Google Drive link processed successfully. {files_count} files have been processed. You can now ask me questions about the content."
                                    print(f"🔊 Announcing Drive processing success: {verbal_msg}")
                                    
                                    # Try fallback_tts first (faster)
                                    try:
                                        result = fallback_tts(verbal_msg)
                                        if result:
                                            print("✅ Drive processing announcement completed via fallback_tts")
                                        else:
                                            print("⚠️ fallback_tts returned False, trying TextToSpeech")
                                            TextToSpeech(verbal_msg)
                                    except Exception as tts_error:
                                        print(f"⚠️ fallback_tts error: {tts_error}, trying TextToSpeech")
                                        try:
                                            TextToSpeech(verbal_msg)
                                            print("✅ Drive processing announcement completed via TextToSpeech")
                                        except Exception as wrapper_error:
                                            print(f"❌ Both TTS methods failed: {wrapper_error}")
                                except Exception as e:
                                    print(f"❌ Error in verbal announcement: {e}")
                                    import traceback
                                    traceback.print_exc()
                            
                            # Start announcement in background thread (non-daemon so it completes)
                            import threading
                            announcement_thread = threading.Thread(target=announce_success_verbally, daemon=False)
                            announcement_thread.start()
                            print("🔊 Started Drive processing announcement thread")
                            
                            # IMPORTANT: Clear the message completely and return - don't send anything to chatbot
                            message = ""
                            message_clean = ""
                            print("Drive link processed successfully - returning without sending to chatbot")
                            return  # Exit early - don't process any remaining text
                        else:
                            # Better error reporting
                            if not result:
                                error_msg = "❌ Error: No result returned from Drive processing"
                            elif not isinstance(result, dict):
                                error_msg = f"❌ Error: Unexpected result type: {type(result)}"
                            else:
                                error_msg = f"❌ Error processing Drive link: {result.get('error', 'Unknown error')}"
                                if result.get('suggestion'):
                                    error_msg += f"\n💡 {result.get('suggestion')}"
                            print(f"Error message: {error_msg}")
                            self.addMessage(error_msg, "Red")
                            self.label.setText("Available...")
                    except Exception as e:
                        import traceback
                        error_details = traceback.format_exc()
                        print(f"Error processing Drive link: {e}")
                        print(f"Full error traceback:\n{error_details}")
                        error_msg = f"❌ Error processing Drive link: {str(e)}"
                        if "No module named" in str(e) or "ImportError" in str(e):
                            error_msg += "\n💡 Missing dependency. Please install: pip install beautifulsoup4 gdown"
                        self.addMessage(error_msg, "Red")
                        self.label.setText("Available...")
                    
                    # Remove the Drive link from message to continue normal processing
                    message = message.replace(drive_link, "").strip()
                    message_clean = message_clean.replace(drive_link, "").strip()
                    
                    # Also remove common processing-related words that might confuse the chatbot
                    processing_words = [
                        "process", "process this", "process:", "process this link", 
                        "process this drive link", "process the link", "process link",
                        "process drive", "process drive link", "process the drive link"
                    ]
                    for word in processing_words:
                        # Case-insensitive replacement
                        import re
                        message = re.sub(re.escape(word), "", message, flags=re.IGNORECASE).strip()
                        message_clean = re.sub(re.escape(word), "", message_clean, flags=re.IGNORECASE).strip()
                    
                    # Clean up any extra whitespace or punctuation
                    message = re.sub(r'\s+', ' ', message).strip()
                    message_clean = re.sub(r'\s+', ' ', message_clean).strip()
                
                # If message is now empty or only contains processing-related words, return
                remaining_text = message.strip().lower()
                if (not message or 
                    not remaining_text or 
                    remaining_text in ["process", "process:", "process this", "process link", "process the link"] or
                    remaining_text.startswith("process")):
                    self.label.setText("Available...")
                    print("Message is empty or only contains processing words - not sending to chatbot")
                    return
            
            # Import the necessary modules
            from Backend.Model import FirstLayerDMM
            from Backend.Chatbot import ChatBot
            from Backend.RealtimeSearchEngine import RealtimeSearchEngine
            from Backend.TextToSpeech import TextToSpeech, interrupt_speech
            from Backend.Automation import Automation
            from asyncio import run
            import subprocess
            import os
            import time
            
            # Get environment variables
            env_vars = dotenv_values(".env")
            Username = env_vars.get("Username", "User")
            Assistantname = env_vars.get("Assistantname", "Assistant")
            
            # Get current mode
            current_mode = self.mode_dropdown.currentText()
            print(f"=" * 50)
            print(f"PROCESSING MESSAGE: '{message}'")
            print(f"CURRENT MODE FROM DROPDOWN: '{current_mode}'")
            print(f"=" * 50)
            
            # Check if user is asking about current mode - handle EARLY before processing
            message_lower = message.lower().strip()
            mode_question_phrases = [
                "which mode", "what mode", "current mode", "in which mode", "what mode are you",
                "which mode are you in", "what mode am i in", "what's the current mode",
                "what mode is active", "which mode is active", "what mode are you currently in",
                "in which more", "which more", "what more", "current more",  # Handle typos
                "in which mode currently", "in which mode are you", "what mode am i",
                "which mode i am", "which mode am i", "tell me which mode", "tell me what mode"
            ]
            
            if any(phrase in message_lower for phrase in mode_question_phrases):
                # Handle mode query directly - return immediately
                print(f"MODE QUESTION DETECTED! Passing mode '{current_mode}' to ChatBot")
                from Backend.Chatbot import ChatBot
                mode_answer = ChatBot(message, mode=current_mode)
                print(f"MODE ANSWER RECEIVED: '{mode_answer[:100]}...'")
                if mode_answer:
                    # Ensure answer mentions JARVIS, not FRIDAY
                    if mode_answer and ("FRIDAY" in mode_answer.upper() or "friday" in mode_answer.lower()):
                        mode_answer = mode_answer.replace("FRIDAY", "JARVIS").replace("friday", "JARVIS").replace("Friday", "JARVIS")
                    
                    from Frontend.GUI import ShowTextToScreen
                    ShowTextToScreen(f"{Assistantname}: {mode_answer}")
                    
                    # Start TTS in thread
                    def speak_mode_in_thread():
                        try:
                            print("=" * 50)
                            print("TTS THREAD STARTED (MODE QUERY)")
                            print("=" * 50)
                            from Backend.TextToSpeech import TextToSpeech, reset_speech_interrupt
                            from Main import check_for_interruption
                            reset_speech_interrupt()
                            def check_interrupt(r=None):
                                from Backend.TextToSpeech import speech_interrupted, interrupt_requested
                                return not (speech_interrupted or interrupt_requested or check_for_interruption())
                            clean_answer = str(mode_answer).strip()
                            print(f"TTS (mode): Original answer: '{clean_answer[:100]}...'")
                            # Remove HTML, markdown, URLs for TTS
                            import re
                            clean_answer = re.sub(r'<[^>]+>', '', clean_answer)
                            clean_answer = re.sub(r'\*\*([^*]+)\*\*', r'\1', clean_answer)
                            clean_answer = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_answer)
                            clean_answer = re.sub(r'http[s]?://[^\s]+', 'link', clean_answer)
                            clean_answer = clean_answer.strip()
                            print(f"TTS (mode): Cleaned answer: '{clean_answer[:100]}...' (length: {len(clean_answer)})")
                            if clean_answer:
                                print("TTS (mode): Calling fallback_tts directly...")
                                try:
                                    from Backend.TextToSpeech import fallback_tts
                                    result = fallback_tts(clean_answer)
                                    print(f"TTS (mode): fallback_tts returned: {result}")
                                    if not result:
                                        TextToSpeech(clean_answer, check_interrupt)
                                    print("TTS (mode): TTS completed")
                                except Exception as e:
                                    print(f"TTS (mode): Error: {e}, trying TextToSpeech wrapper...")
                                    try:
                                        TextToSpeech(clean_answer, check_interrupt)
                                    except:
                                        print(f"TTS (mode): All methods failed")
                            else:
                                print("TTS (mode): ERROR - Cleaned answer is empty!")
                            self.label.setText("Available...")
                        except Exception as e:
                            print(f"TTS (mode): ERROR in thread: {e}")
                            import traceback
                            print(traceback.format_exc())
                    import threading
                    speech_thread = threading.Thread(target=speak_mode_in_thread, daemon=False)  # Changed to non-daemon
                    speech_thread.start()
                    print("TTS (mode): Thread started")
                    self.label.setText("Available...")
                    return
            
            # Show thinking status immediately
            self.label.setText(f"Thinking... ({current_mode})")
            self.update()  # Force immediate GUI update
            
            # Check for simple greetings FIRST - handle them immediately to prevent duplicates
            greeting_words = ["hello", "hi", "hey", "namaste", "hola", "greetings"]
            message_lower_check = message.lower().strip()
            words_check = message_lower_check.split()
            is_simple_greeting = (len(words_check) <= 2 and 
                                 any(word in greeting_words for word in words_check) and
                                 all(word in greeting_words + ["jarvis", "there"] for word in words_check))
            
            if is_simple_greeting:
                # Handle greeting directly - only ONE response
                print(f"SIMPLE GREETING DETECTED: '{message}' - handling directly")
                from Backend.Chatbot import ChatBot
                Answer = ChatBot(message, mode=current_mode)
                if Answer:
                    # Display answer and return immediately - don't process further
                    Answer = Answer.replace("FRIDAY", "JARVIS").replace("friday", "JARVIS").replace("Friday", "JARVIS")
                    self.addMessage(f"{Assistantname}: {Answer}", "LightBlue")
                    # Write to file
                    try:
                        current_content = ""
                        with open(rf'{TempDirPath}\Responses.data', 'r', encoding='utf-8') as f:
                            current_content = f.read()
                        if current_content and not current_content.endswith('\n'):
                            current_content += '\n'
                        current_content += f"{Assistantname}: {Answer}\n"
                        with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as f:
                            f.write(current_content)
                    except:
                        pass
                    # Start TTS
                    def speak_greeting():
                        try:
                            from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt
                            reset_speech_interrupt()
                            fallback_tts(Answer)
                        except:
                            pass
                    import threading
                    threading.Thread(target=speak_greeting, daemon=True).start()
                    # Update status
                    try:
                        from Frontend.GUI import GetMicrophoneStatus
                        mic_status = GetMicrophoneStatus()
                        if mic_status and mic_status.lower() == "true":
                            self.label.setText("Listening...")
                            SetAsssistantStatus("Listening...")
                        else:
                            self.label.setText("Available...")
                            SetAsssistantStatus("Available...")
                    except:
                        self.label.setText("Available...")
                    # Remove from processing set
                    if hasattr(self, '_processing_messages'):
                        self._processing_messages.discard(message_key)
                    return  # EXIT IMMEDIATELY - don't process further
            
            # Process through the decision model
            Decision = FirstLayerDMM(message)
            print(f"Decision: {Decision}")
            
            # Check for different types of queries
            G = any([i for i in Decision if i.startswith("general")])
            R = any([i for i in Decision if i.startswith("realtime")])
            
            # Handle automation tasks - check ALL decision items for task commands
            functions = ["open", "close", "play", "system", "content", "google search", "youtube search", "spotify", "send email", "send whatsapp", "schedule meeting", "sync crm"]
            TaskExecution = False
            
            # Check if ANY decision item contains a task command
            for queries in Decision:
                if not TaskExecution:
                    # Check if query starts with any function OR contains it (for cases like "general open youtube")
                    query_lower = queries.lower()
                    if any(query_lower.startswith(func) or f" {func} " in f" {query_lower} " for func in functions):
                        print(f"TASK DETECTED: {queries}")
                        self.label.setText("Executing task...")
                        self.update()  # Force immediate GUI update
                        try:
                            # Execute automation with the full Decision list
                            from asyncio import run
                            from Backend.Automation import Automation
                            run(Automation(list(Decision)))
                            self.label.setText("Task completed...")
                            self.addMessage(f"{Assistantname}: Task completed successfully!", "Green")
                            self.update()  # Force immediate GUI update
                            
                            # Speak the success message
                            try:
                                from Backend.TextToSpeech import fallback_tts
                                fallback_tts("Task completed successfully!")
                            except:
                                pass
                            
                            # Minimal delay for faster response
                            time.sleep(0.1)
                            # Check mic status and set to Listening if mic is ON
                            try:
                                from Frontend.GUI import GetMicrophoneStatus
                                mic_status = GetMicrophoneStatus()
                                if mic_status and mic_status.lower() == "true":
                                    self.label.setText("Listening...")
                                    SetAsssistantStatus("Listening...")
                                else:
                                    self.label.setText("Available...")
                                    SetAsssistantStatus("Available...")
                            except:
                                self.label.setText("Available...")
                            return
                        except Exception as e:
                            print(f"Automation error: {e}")
                            import traceback
                            print(traceback.format_exc())
                            self.label.setText("Task failed...")
                            self.addMessage(f"{Assistantname}: Task failed - {str(e)}", "Red")
                            self.update()  # Force immediate GUI update
                            
                            # Speak the error message
                            try:
                                from Backend.TextToSpeech import fallback_tts
                                fallback_tts(f"Task failed: {str(e)[:50]}")
                            except:
                                pass
                            
                            time.sleep(0.1)  # Reduced delay for faster response
                            # Check mic status and set to Listening if mic is ON
                            try:
                                from Frontend.GUI import GetMicrophoneStatus
                                mic_status = GetMicrophoneStatus()
                                if mic_status and mic_status.lower() == "true":
                                    self.label.setText("Listening...")
                                    SetAsssistantStatus("Listening...")
                                else:
                                    self.label.setText("Available...")
                                    SetAsssistantStatus("Available...")
                            except:
                                self.label.setText("Available...")
                            return
            
            # Ensure real-time questions use Groq API
            R = any([i for i in Decision if i.startswith("realtime")])
            G = any([i for i in Decision if i.startswith("general")])
            
            # CRITICAL: Check for "who is" questions FIRST - these MUST go to realtime search
            # Don't process through general path if it's a real-time question
            message_lower = message.lower().strip()
            realtime_keywords = [
                "who is the", "who is", "current president", "current prime minister", 
                "president of", "prime minister of", "current", "now", "today", 
                "recent", "latest", "news", "headline", "what's happening",
                "ceo of", "leader of", "ruler of", "king of", "queen of",
                "gdp", "gross domestic product", "ranking", "rank", "position", 
                "world ranking", "terms of", "in terms of", "economic", "economy",
                "largest economy", "richest country", "population", "current population",
                "unemployment rate", "inflation rate", "stock market", "cryptocurrency",
                "bitcoin price", "exchange rate", "currency", "market cap"
            ]
            
            # If it's a real-time question, FORCE it to realtime path (don't use general)
            if any(keyword in message_lower for keyword in realtime_keywords):
                if not R:
                    print(f"FORCING real-time question (was misclassified as general): {message}")
                    R = True
                    G = False  # Don't process through general path
            
            if R:
                # Real-time questions - always use Groq API via RealtimeSearchEngine
                # DON'T process through general path - only use realtime
                print(f"REALTIME QUESTION - Passing mode '{current_mode}' to RealtimeSearchEngine")
                # Update status IMMEDIATELY after user message is displayed
                self.label.setText("Searching...")
                self.update()
                QApplication.processEvents()  # Force GUI update to show status
                try:
                    # Use the original message directly (don't merge with general)
                    # Pass mode context to RealtimeSearchEngine for mode-aware responses
                    mode_context_query = f"[Mode: {current_mode}] {message}"
                    Answer = RealtimeSearchEngine(mode_context_query, mode=current_mode)
                    
                    # Verify we got a valid answer
                    if not Answer or Answer.strip() == "":
                        print("WARNING: RealtimeSearchEngine returned empty answer, trying fallback")
                        # Only try fallback if realtime completely failed
                        Answer = self.process_with_mode(message, Decision, current_mode, Assistantname)
                    else:
                        print(f"RealtimeSearchEngine returned answer (length: {len(Answer)})")
                        # Filter out "I don't know" type responses - if answer contains negative phrases, try again
                        negative_phrases = ["couldn't find", "don't know", "no information", "notable", "well-known", "private individual", "limited information"]
                        if any(phrase in Answer.lower() for phrase in negative_phrases) and len(Answer) < 200:
                            print("WARNING: Answer contains negative phrases, trying to get better answer")
                            # Try one more time with clearer prompt
                            try:
                                better_answer = RealtimeSearchEngine(f"[Mode: {current_mode}] User is asking: {message}. Use search results to provide detailed information.", mode=current_mode)
                                if better_answer and len(better_answer) > 100:
                                    Answer = better_answer
                                    print("Got better answer on retry")
                            except:
                                pass
                except Exception as e:
                    print(f"RealtimeSearchEngine error: {e}")
                    import traceback
                    print(traceback.format_exc())
                    # Only fallback if realtime completely failed
                    Answer = self.process_with_mode(message, Decision, current_mode, Assistantname)
            else:
                # Process with mode-specific context (handles general questions and all other types)
                print(f"GENERAL QUESTION - Passing mode '{current_mode}' to process_with_mode")
                # Update status IMMEDIATELY after user message is displayed
                self.label.setText("Thinking...")
                self.update()
                QApplication.processEvents()  # Force GUI update to show status
                Answer = self.process_with_mode(message, Decision, current_mode, Assistantname)
            
            # Ensure we always have an answer (but only if we don't have one already)
            if not Answer or Answer.strip() == "":
                # Fallback: try ChatBot directly with mode
                try:
                    print("No answer received, trying ChatBot directly...")
                    Answer = ChatBot(message, mode=current_mode)
                except Exception as e:
                    print(f"ChatBot fallback error: {e}")
                    Answer = "I'm here to help! Please try asking your question again or rephrase it."
            
            # CRITICAL: Only display ONE response - check if we already displayed one
            if Answer and Answer.strip():
                # Normalize answer for comparison (remove extra whitespace)
                answer_normalized = ' '.join(Answer.strip().split())
                
                # Check if this exact answer was already displayed recently (increased window to 5 seconds)
                if hasattr(self, '_last_displayed_answer'):
                    import time
                    current_time = time.time()
                    # If same answer was displayed within last 5 seconds, skip
                    if (self._last_displayed_answer[0] == answer_normalized and 
                        (current_time - self._last_displayed_answer[1]) < 5.0):
                        print(f"Duplicate answer detected, skipping display: '{Answer[:50]}...'")
                        # Remove from processing set
                        self._processing_messages.discard(message_key)
                        return
                
                # Mark this answer as displayed
                import time
                self._last_displayed_answer = (answer_normalized, time.time())
            
            if Answer:
                # Ensure answer mentions JARVIS, not FRIDAY - comprehensive replacement
                if Answer:
                    Answer = Answer.replace("FRIDAY", "JARVIS").replace("friday", "JARVIS").replace("Friday", "JARVIS")
                    # Remove any corrections about name
                    Answer = Answer.replace("but I'm actually FRIDAY", "I'm JARVIS").replace("but I'm actually friday", "I'm JARVIS")
                    Answer = Answer.replace("actually FRIDAY", "JARVIS").replace("actually friday", "JARVIS")
                
                # Display answer in chat immediately (user message was already added in send_text_message)
                self.addMessage(f"{Assistantname}: {Answer}", "LightBlue")
                
                # Also append answer to Responses.data file to maintain order (append, don't overwrite)
                try:
                    current_content = ""
                    with open(rf'{TempDirPath}\Responses.data', 'r', encoding='utf-8') as f:
                        current_content = f.read()
                    # Append answer to file content
                    if current_content and not current_content.endswith('\n'):
                        current_content += '\n'
                    current_content += f"{Assistantname}: {Answer}\n"
                    with open(rf'{TempDirPath}\Responses.data', 'w', encoding='utf-8') as f:
                        f.write(current_content)
                except Exception as e:
                    print(f"Error writing answer to file: {e}")
                
                self.label.setText("Answering...")
                self.update()  # Force immediate GUI update
                
                # Start TTS immediately in a separate thread for faster response (NO DELAYS)
                def speak_in_thread():
                    try:
                        # Check if this exact text was already spoken recently (within 5 seconds)
                        from Main import last_spoken_text, last_spoken_time
                        import time
                        answer_normalized_tts = ' '.join(Answer.strip().split())
                        current_tts_time = time.time()
                        
                        if (last_spoken_text == answer_normalized_tts and 
                            (current_tts_time - last_spoken_time) < 5.0):
                            print(f"GUI TTS: Duplicate speech detected, skipping TTS: '{Answer[:50]}...'")
                            return
                        
                        # Mark this as about to be spoken (update global variables)
                        from Main import last_spoken_text, last_spoken_time
                        import Main
                        Main.last_spoken_text = answer_normalized_tts
                        Main.last_spoken_time = current_tts_time
                        
                        print("=" * 50)
                        print("TTS THREAD STARTED")
                        print("=" * 50)
                        
                        # Import TTS functions
                        from Backend.TextToSpeech import TextToSpeech, speech_interrupted, interrupt_requested, reset_speech_interrupt
                        from Main import check_for_interruption
                        
                        # Reset interrupt flags before starting new speech
                        reset_speech_interrupt()
                        
                        # NO DELAY - start speaking immediately for fastest response
                        
                        def check_interrupt(r=None):
                            # Simplified interrupt check - only check if explicitly interrupted
                            from Backend.TextToSpeech import speech_interrupted, interrupt_requested
                            should_interrupt = (speech_interrupted or interrupt_requested)
                            if should_interrupt:
                                print(f"TTS interrupt check: speech_interrupted={speech_interrupted}, interrupt_requested={interrupt_requested}")
                            return not should_interrupt
                        
                        # Clean answer text for TTS (remove HTML, markdown, but keep content)
                        clean_answer = str(Answer).strip()
                        print(f"Original answer length: {len(clean_answer)}")
                        print(f"Original answer preview: {clean_answer[:150]}...")
                        
                        # Remove markdown links but keep the link text
                        import re
                        clean_answer = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_answer)
                        # Remove URLs but keep readable text (replace with "link" or remove)
                        clean_answer = re.sub(r'https?://[^\s]+', 'link', clean_answer)
                        clean_answer = re.sub(r'www\.[^\s]+', 'link', clean_answer)
                        # Remove any HTML tags but keep content
                        clean_answer = re.sub(r'<[^>]+>', '', clean_answer)
                        # Remove style attributes that might be in text
                        clean_answer = re.sub(r'\.?\s*style=["\'][^"\']*["\']', '', clean_answer)
                        # Remove HTML entities
                        clean_answer = clean_answer.replace('&nbsp;', ' ').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
                        # Remove extra whitespace
                        clean_answer = ' '.join(clean_answer.split())
                        
                        # Ensure we have text to speak
                        if not clean_answer or not clean_answer.strip():
                            print("WARNING: Cleaning removed all text, using minimally cleaned original")
                            # If cleaning removed everything, use original (but clean it minimally)
                            clean_answer = str(Answer).strip()
                            clean_answer = re.sub(r'<[^>]+>', '', clean_answer)
                            clean_answer = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_answer)
                            clean_answer = ' '.join(clean_answer.split())
                        
                        # Always try to speak, even if cleaning removed some text
                        if not clean_answer or not clean_answer.strip():
                            print("WARNING: Cleaned answer is empty, using original answer for TTS")
                            clean_answer = str(Answer).strip()
                            # Minimal cleaning - just remove HTML and markdown
                            clean_answer = re.sub(r'<[^>]+>', '', clean_answer)
                            clean_answer = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_answer)
                            clean_answer = ' '.join(clean_answer.split())
                        
                        if clean_answer and clean_answer.strip():
                            print(f"TTS: Final cleaned answer length: {len(clean_answer)}")
                            print(f"TTS: Final cleaned answer: '{clean_answer[:200]}...'")
                            print("TTS: Calling TextToSpeech function now...")
                            
                            # Use direct fallback_tts for maximum reliability - ALWAYS CALL IT
                            try:
                                print(f"TTS: Calling fallback_tts directly with text: '{clean_answer[:100]}...'")
                                from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt
                                # Ensure interrupts are reset
                                reset_speech_interrupt()
                                # Call fallback_tts directly - it's more reliable
                                direct_result = fallback_tts(clean_answer)
                                print(f"TTS: fallback_tts returned: {direct_result}")
                                if direct_result:
                                    print("TTS: Speech completed successfully via fallback_tts")
                                else:
                                    print("WARNING: fallback_tts returned False - trying TextToSpeech wrapper")
                                    # Fallback to wrapper if direct call fails
                                    result = TextToSpeech(clean_answer, check_interrupt)
                                    print(f"TTS: TextToSpeech wrapper returned: {result}")
                            except Exception as tts_error:
                                print(f"ERROR in TTS: {tts_error}")
                                import traceback
                                print(traceback.format_exc())
                                # Last resort - try the wrapper
                                try:
                                    print("TTS: Trying TextToSpeech wrapper as last resort...")
                                    result = TextToSpeech(clean_answer, check_interrupt)
                                    print(f"TTS: TextToSpeech wrapper returned: {result}")
                                except Exception as wrapper_error:
                                    print(f"TTS: All TTS methods failed: {wrapper_error}")
                                    # Final fallback - try simple pyttsx3
                                    try:
                                        import pyttsx3
                                        engine = pyttsx3.init()
                                        engine.setProperty('rate', 200)
                                        engine.say(clean_answer[:500])  # Limit length
                                        engine.runAndWait()
                                        print("TTS: Simple pyttsx3 fallback succeeded")
                                    except Exception as simple_error:
                                        print(f"TTS: All TTS methods including pyttsx3 failed: {simple_error}")
                            
                            print("TTS: TextToSpeech function call finished")
                        else:
                            print("ERROR: Answer is completely empty after cleaning, cannot speak")
                            print(f"Original answer was: {Answer[:200] if Answer else 'None'}...")
                            
                    except Exception as e:
                        import traceback
                        print(f"TTS THREAD ERROR: {e}")
                        print(f"TTS THREAD TRACEBACK: {traceback.format_exc()}")
                    finally:
                        try:
                            # Check mic status and set to Listening if mic is ON
                            from Frontend.GUI import GetMicrophoneStatus
                            mic_status = GetMicrophoneStatus()
                            if mic_status and mic_status.lower() == "true":
                                self.label.setText("Listening...")
                                SetAsssistantStatus("Listening...")
                                print("TTS THREAD: Status set to Listening... (mic is ON)")
                            else:
                                self.label.setText("Available...")
                                SetAsssistantStatus("Available...")
                                print("TTS THREAD: Status set to Available... (mic is OFF)")
                        except Exception as e:
                            print(f"Error setting status after TTS: {e}")
                            try:
                                self.label.setText("Available...")
                            except:
                                pass
                        print("=" * 50)
                        print("TTS THREAD FINISHED")
                        print("=" * 50)
                
                # Start speaking immediately - use direct call for reliability
                import threading
                
                # CRITICAL: Call TTS directly first (synchronous) to ensure it works
                try:
                    print(f"TTS: Calling fallback_tts DIRECTLY (synchronous) for immediate speech")
                    from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt
                    reset_speech_interrupt()
                    
                    # Clean answer for TTS
                    import re
                    clean_for_tts = str(Answer).strip()
                    clean_for_tts = re.sub(r'<[^>]+>', '', clean_for_tts)
                    clean_for_tts = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_for_tts)
                    clean_for_tts = re.sub(r'https?://[^\s]+', 'link', clean_for_tts)
                    clean_for_tts = ' '.join(clean_for_tts.split())
                    
                    # Call TTS in a thread but ensure it actually runs
                    def direct_tts_call():
                        try:
                            print(f"TTS DIRECT: Speaking '{clean_for_tts[:100]}...'")
                            result = fallback_tts(clean_for_tts)
                            print(f"TTS DIRECT: Result = {result}")
                            if not result:
                                print("TTS DIRECT: fallback_tts returned False, trying simple pyttsx3")
                                import pyttsx3
                                engine = pyttsx3.init()
                                engine.setProperty('rate', 200)
                                engine.setProperty('volume', 1.0)
                                engine.say(clean_for_tts[:500])
                                engine.runAndWait()
                                print("TTS DIRECT: Simple pyttsx3 succeeded")
                        except Exception as e:
                            print(f"TTS DIRECT ERROR: {e}")
                            import traceback
                            print(traceback.format_exc())
                    
                    # Start in thread (non-blocking)
                    tts_thread = threading.Thread(target=direct_tts_call, daemon=False)
                    tts_thread.start()
                    print(f"TTS thread started - speech should begin immediately (Answer length: {len(Answer)})")
                    
                    # Also start the full TTS thread as backup
                    speech_thread = threading.Thread(target=speak_in_thread, daemon=False)
                    speech_thread.start()
                    print("Full TTS thread also started as backup")
                except Exception as tts_start_error:
                    print(f"TTS start error: {tts_start_error}")
                    import traceback
                    print(traceback.format_exc())
                    # Try the original thread method
                    speech_thread = threading.Thread(target=speak_in_thread, daemon=False)
                    speech_thread.start()
            else:
                # No answer received - check mic status and set to Listening if mic is ON
                try:
                    from Frontend.GUI import GetMicrophoneStatus
                    mic_status = GetMicrophoneStatus()
                    if mic_status and mic_status.lower() == "true":
                        self.label.setText("Listening...")
                        SetAsssistantStatus("Listening...")
                    else:
                        self.label.setText("Available...")
                        SetAsssistantStatus("Available...")
                except:
                    self.label.setText("Available...")
            
        except Exception as e:
            print(f"Error processing text message: {e}")
            import traceback
            print(traceback.format_exc())
            self.addMessage(f"Error: {str(e)}", "Red")
            self.label.setText("Error occurred...")
            # Check mic status and set to Listening if mic is ON
            time.sleep(1)
            try:
                from Frontend.GUI import GetMicrophoneStatus
                mic_status = GetMicrophoneStatus()
                if mic_status and mic_status.lower() == "true":
                    self.label.setText("Listening...")
                    SetAsssistantStatus("Listening...")
                else:
                    self.label.setText("Available...")
                    SetAsssistantStatus("Available...")
            except:
                self.label.setText("Available...")
        finally:
            # Always remove from processing set when done
            if hasattr(self, '_processing_messages'):
                self._processing_messages.discard(message_key)
    
    def upload_file(self):
        """Handle file upload"""
        try:
            from PyQt5.QtWidgets import QFileDialog
            file_path, _ = QFileDialog.getOpenFileName(
                self, 
                "Select File to Upload", 
                "", 
                "All Files (*);;Text Files (*.txt *.md *.csv *.log);;PDF Files (*.pdf);;Image Files (*.png *.jpg *.jpeg *.gif *.bmp *.tiff);;Word Documents (*.docx);;Excel Files (*.xlsx *.xls);;JSON Files (*.json);;XML Files (*.xml)"
            )
            
            if file_path:
                # Extract filename from path
                import os
                filename = os.path.basename(file_path)
                
                # Add file upload message to chat
                self.addMessage(f"📁 Uploaded file: {filename}", "LightBlue")
                
                # Process the file based on its type
                self.process_uploaded_file(file_path, filename)
                
        except Exception as e:
            print(f"Error uploading file: {e}")
            self.addMessage(f"Error uploading file: {str(e)}", "Red")
    
    def process_uploaded_file(self, file_path, filename):
        """Process uploaded file and provide analysis"""
        try:
            import os
            file_extension = os.path.splitext(filename)[1].lower()
            
            # Show processing status
            self.label.setText("Processing file...")
            self.update()
            
            content = ""
            file_type = "Unknown"
            
            if file_extension in ['.txt', '.md', '.csv', '.log']:
                # Process text files
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    file_type = "Text File"
                except UnicodeDecodeError:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        content = f.read()
                    file_type = "Text File (Legacy Encoding)"
                
            elif file_extension in ['.pdf']:
                # Process PDF files
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        content = ""
                        for page_num in range(min(5, len(pdf_reader.pages))):  # First 5 pages
                            page = pdf_reader.pages[page_num]
                            content += page.extract_text() + "\n"
                    file_type = "PDF Document"
                except Exception as e:
                    content = f"Error reading PDF: {str(e)}"
                    file_type = "PDF Document (Error)"
                
            elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
                # Process image files with OCR
                try:
                    import pytesseract
                    from PIL import Image
                    import cv2
                    import numpy as np
                    
                    # Load image
                    image = cv2.imread(file_path)
                    
                    if image is None:
                        content = "Could not read image file"
                        file_type = "Image (Error)"
                    else:
                        # Convert to RGB
                        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
                        
                        # Extract text using OCR
                        content = pytesseract.image_to_string(image_rgb)
                        
                        if not content.strip():
                            content = "No text detected in image. This might be a photo without readable text."
                        
                        file_type = "Image with Text"
                except Exception as e:
                    content = f"Error processing image: {str(e)}"
                    file_type = "Image (OCR Error)"
                
            elif file_extension in ['.docx']:
                # Process Word documents
                try:
                    from docx import Document
                    doc = Document(file_path)
                    content = ""
                    for paragraph in doc.paragraphs:
                        content += paragraph.text + "\n"
                    file_type = "Word Document"
                except Exception as e:
                    content = f"Error reading Word document: {str(e)}"
                    file_type = "Word Document (Error)"
                
            elif file_extension in ['.xlsx', '.xls']:
                # Process Excel files
                try:
                    import pandas as pd
                    excel_file = pd.ExcelFile(file_path)
                    content = "Excel File Contents:\n\n"
                    
                    for sheet_name in excel_file.sheet_names[:3]:  # First 3 sheets
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        content += f"Sheet: {sheet_name}\n"
                        content += df.to_string(max_rows=10, max_cols=10) + "\n\n"
                    
                    file_type = "Excel Spreadsheet"
                except Exception as e:
                    content = f"Error reading Excel file: {str(e)}"
                    file_type = "Excel File (Error)"
                
            elif file_extension in ['.ppt', '.pptx']:
                # Process PowerPoint files
                try:
                    from pptx import Presentation  # type: ignore[reportMissingImports]  # Optional dependency
                    prs = Presentation(file_path)
                    content = f"PowerPoint Presentation: {filename}\n\n"
                    
                    slide_count = 0
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_count += 1
                        content += f"Slide {slide_num}:\n"
                        
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                content += shape.text + "\n"
                            # Also check for tables
                            if hasattr(shape, "table"):
                                table_text = ""
                                for row in shape.table.rows:
                                    row_text = " | ".join([cell.text for cell in row.cells])
                                    table_text += row_text + "\n"
                                if table_text:
                                    content += "Table:\n" + table_text
                        
                        content += "\n"
                    
                    content += f"\nTotal slides: {slide_count}"
                    file_type = "PowerPoint Presentation"
                except ImportError:
                    content = f"PowerPoint file uploaded: {filename}\n\nTo process PPT files, please install python-pptx: pip install python-pptx"
                    file_type = "PowerPoint (Library Not Installed)"
                except Exception as e:
                    content = f"Error reading PowerPoint file: {str(e)}"
                    file_type = "PowerPoint (Error)"
                
            elif file_extension in ['.json']:
                # Process JSON files
                try:
                    import json
                    with open(file_path, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                    content = json.dumps(data, indent=2)
                    file_type = "JSON File"
                except Exception as e:
                    content = f"Error reading JSON file: {str(e)}"
                    file_type = "JSON File (Error)"
                
            elif file_extension in ['.xml']:
                # Process XML files
                try:
                    from bs4 import BeautifulSoup
                    with open(file_path, 'r', encoding='utf-8') as f:
                        soup = BeautifulSoup(f, 'xml')
                    content = soup.prettify()
                    file_type = "XML File"
                except Exception as e:
                    content = f"Error reading XML file: {str(e)}"
                    file_type = "XML File (Error)"
                
            else:
                # Generic file upload
                content = f"File uploaded: {filename}\nFile type: {file_extension}\nFile size: {os.path.getsize(file_path)} bytes"
                file_type = f"File ({file_extension})"
            
            # Display file content
            if content:
                preview = content[:1000] + "..." if len(content) > 1000 else content
                self.addMessage(f"📄 {file_type}: {filename}\n\nContent Preview:\n{preview}", "White")
            else:
                self.addMessage(f"📄 {file_type}: {filename}\n\nNo readable content found.", "White")
            
            # Process file for sales knowledge if it's a document
            try:
                from Backend.DocumentProcessor import process_document
                process_result = process_document(file_path, filename)
                if process_result.get("success"):
                    self.addMessage(f"✅ File processed and stored in sales knowledge base! ({process_result.get('entries_created', 0)} entries created)", "LightGreen")
                else:
                    self.addMessage(f"Note: {process_result.get('error', 'Could not process for sales knowledge')}", "Yellow")
            except Exception as e:
                print(f"Error processing file for sales knowledge: {e}")
            
            # Process file content with Groq API for understanding
            try:
                from Backend.RealtimeSearchEngine import RealtimeSearchEngine
                # Use Groq API to understand the file content
                # Limit content to first 2000 chars for API
                file_content_preview = content[:2000] if len(content) > 2000 else content
                file_understanding_prompt = f"""I have uploaded a file named '{filename}' (Type: {file_type}). 
The file content is:
{file_content_preview}

Please analyze and summarize this file. What is it about? What are the key points or important information in it? 
If it's a document, resume, or data file, provide a comprehensive summary."""
                
                self.label.setText("Understanding file...")
                self.update()
                
                # Use Groq API to understand the file
                file_analysis = RealtimeSearchEngine(file_understanding_prompt)
                self.addMessage(f"JARVIS: {file_analysis}", "White")
                
            except Exception as e:
                print(f"Error analyzing file with Groq API: {e}")
                # Fallback message
                self.addMessage("JARVIS: File uploaded successfully. I can answer questions about this file. What would you like to know?", "LightGreen")
            
            self.label.setText("Available...")
            
        except Exception as e:
            print(f"Error processing file: {e}")
            self.addMessage(f"Error processing file: {str(e)}", "Red")
            self.label.setText("Error occurred...")
    
    def on_mode_changed(self, mode):
        """Handle mode selection change - updates ModeManager and plays TTS confirmation"""
        try:
            from Backend.ModeManager import set_mode, get_mode_manager, get_current_mode
            from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt
            
            # Try to import logger, but don't fail if it doesn't work
            try:
                from Backend.Logger import log_mode_change, log_tts
                logger_available = True
            except ImportError:
                logger_available = False
                def log_mode_change(old, new):
                    pass
                def log_tts(text, status):
                    pass
            
            old_mode = get_current_mode()
            
            # Update mode in ModeManager
            if set_mode(mode, source="gui"):
                if logger_available:
                    try:
                        log_mode_change(old_mode, mode)
                    except:
                        pass
                
                # Show mode change notification
                self.addMessage(f"🔄 Switched to {mode}", "LightBlue")
                
                # Get mode-specific guidance
                mode_manager = get_mode_manager()
                mode_guidance = mode_manager.get_mode_guidance(mode)
                
                # Show guidance in chat
                self.addMessage(f"JARVIS: {mode_guidance}", "White")
                
                # CRITICAL: Reset status to "Listening..." if mic is ON, otherwise "Available..."
                # This ensures voice recognition continues working after mode switch
                try:
                    from Frontend.GUI import GetMicrophoneStatus
                    mic_status = GetMicrophoneStatus()
                    if mic_status and mic_status.lower() == "true":
                        self.label.setText("Listening...")
                        from Frontend.GUI import SetAsssistantStatus
                        SetAsssistantStatus("Listening...")
                        print(f"Mode switch: Mic is ON, status set to 'Listening...'")
                    else:
                        self.label.setText("Available...")
                        from Frontend.GUI import SetAsssistantStatus
                        SetAsssistantStatus("Available...")
                        print(f"Mode switch: Mic is OFF, status set to 'Available...'")
                except Exception as status_error:
                    print(f"Error setting status after mode switch: {status_error}")
                    # Fallback: just set to listening to be safe
                    self.label.setText("Listening...")
                
                # Play TTS confirmation verbally - use threading to avoid blocking GUI
                def announce_mode_verbally():
                    try:
                        import time
                        # Reset interrupt flags FIRST before any delay
                        reset_speech_interrupt()
                        
                        # Small delay to ensure GUI updates first
                        time.sleep(0.2)
                        
                        # Create shorter, clearer announcement
                        short_announcement = f"Now in {mode} mode."
                        print(f"=== MODE ANNOUNCEMENT: {short_announcement} ===")
                        
                        # Ensure interrupt flags are still reset
                        from Backend.TextToSpeech import speech_interrupted, interrupt_requested
                        if speech_interrupted or interrupt_requested:
                            print("=== Resetting interrupt flags again ===")
                            reset_speech_interrupt()
                        
                        # Try fallback_tts first (faster and more reliable)
                        try:
                            print(f"=== MODE SWITCH TTS: Calling fallback_tts with: '{short_announcement}' ===")
                            # Force reset interrupt flags before TTS
                            reset_speech_interrupt()
                            result = fallback_tts(short_announcement)
                            print(f"=== MODE SWITCH TTS: fallback_tts returned: {result} ===")
                            if result:
                                print("=== MODE SWITCH TTS: TTS announcement successful! ===")
                                if logger_available:
                                    try:
                                        log_tts(short_announcement, "start")
                                    except:
                                        pass
                                return
                            else:
                                print("=== MODE SWITCH TTS: fallback_tts returned False, trying simple pyttsx3 ===")
                        except Exception as fallback_error:
                            print(f"=== MODE SWITCH TTS ERROR: {fallback_error} ===")
                            import traceback
                            traceback.print_exc()
                        
                        # If fallback fails, try simple pyttsx3 directly
                        try:
                            print(f"=== Trying simple pyttsx3 directly ===")
                            import pyttsx3
                            simple_engine = pyttsx3.init()
                            simple_engine.setProperty('rate', 200)
                            simple_engine.setProperty('volume', 1.0)
                            # Clean text
                            clean_text = short_announcement.replace('\n', ' ').strip()
                            if clean_text:
                                simple_engine.say(clean_text)
                                simple_engine.runAndWait()
                                print(f"=== Simple pyttsx3 succeeded ===")
                                return
                        except Exception as simple_error:
                            print(f"=== Simple pyttsx3 error: {simple_error} ===")
                            import traceback
                            traceback.print_exc()
                        
                        # Last resort: try TextToSpeech wrapper
                        try:
                            from Backend.TextToSpeech import TextToSpeech
                            def dummy_check():
                                return True
                            print(f"=== Trying TextToSpeech wrapper as last resort ===")
                            TextToSpeech(short_announcement, dummy_check)
                            print(f"=== Used TextToSpeech wrapper ===")
                        except Exception as wrapper_error:
                            print(f"=== TextToSpeech wrapper error: {wrapper_error} ===")
                            import traceback
                            traceback.print_exc()
                        
                    except Exception as tts_error:
                        print(f"=== CRITICAL Error in announce_mode_verbally: {tts_error} ===")
                        import traceback
                        traceback.print_exc()
                
                # Run TTS in a separate thread to avoid blocking GUI
                # Use NON-DAEMON thread to ensure it completes
                import threading
                tts_thread = threading.Thread(target=announce_mode_verbally, daemon=False)
                tts_thread.start()
                print(f"=== TTS thread started for mode announcement (non-daemon) ===")
                
                # Don't wait for thread, but ensure it's running
                import time
                time.sleep(0.1)  # Small delay to let thread start
                if not tts_thread.is_alive():
                    print(f"=== WARNING: TTS thread died immediately, trying direct call ===")
                    # Try direct call as fallback
                    try:
                        from Backend.TextToSpeech import fallback_tts, reset_speech_interrupt
                        reset_speech_interrupt()
                        fallback_tts(f"Now in {mode} mode.")
                    except Exception as direct_error:
                        print(f"=== Direct TTS call also failed: {direct_error} ===")
                
            else:
                self.addMessage(f"Error: Could not switch to {mode}", "Red")
                
        except Exception as e:
            print(f"Error changing mode: {e}")
            import traceback
            traceback.print_exc()
            # Don't show technical errors to user, just log them
            # Only show user-friendly messages
            # Never show technical errors to user - they're handled internally
            # Only show user-friendly messages for actual mode switching failures
            error_str = str(e).lower()
            if "log_mode_change" not in error_str and "logger" not in error_str and "name" not in error_str:
                # Only show if it's a real mode switching error, not a logging error
                pass  # Don't show errors - mode switching will work anyway
    
    def process_with_mode(self, message, Decision, mode, Assistantname):
        """Process message with mode-specific context"""
        try:
            from Backend.Chatbot import ChatBot
            from Backend.RealtimeSearchEngine import RealtimeSearchEngine
            import time
            
            # Mode-specific system prompts with sales-focused modes
            mode_prompts = {
                "General Assistant": "You are JARVIS, a helpful AI assistant ready to assist with any questions or tasks.",
                
                "Sales Assistant": "You are JARVIS in Sales Assistant Mode. Help with sales activities including lead management, pitch generation, follow-ups, and sales strategies. Use stored sales knowledge (documents, leads, products) to provide personalized, actionable sales advice. Be friendly, persuasive, professional, and focused on closing deals."
            }
            
            # Get mode-specific prompt
            mode_prompt = mode_prompts.get(mode, mode_prompts["General Assistant"])
            
            # Check for different types of queries
            G = any([i for i in Decision if i.startswith("general")])
            R = any([i for i in Decision if i.startswith("realtime")])
            
            Answer = None
            
            # Handle realtime queries with mode context
            if R:
                # Status already set in process_text_message, but ensure it's visible
                if self.label.text() != "Searching...":
                    self.label.setText("Searching...")
                    self.update()
                    QApplication.processEvents()
                
                try:
                    # Add mode context to the query
                    mode_context = f"[Mode: {mode}] {mode_prompt}\n\nUser Query: {message}"
                    
                    if G and R:
                        Merged_query = " and ".join(
                            [" ".join(i.split()[1:]) for i in Decision if i.startswith("general") or i.startswith("realtime")]
                        )
                        # Pass mode to RealtimeSearchEngine
                        Answer = RealtimeSearchEngine(f"{mode_context}\n\nQuery: {Merged_query}", mode=mode)
                    else:
                        for queries in Decision:
                            if "realtime" in queries:
                                QueryFinal = queries.replace("realtime", "")
                                # Pass mode to RealtimeSearchEngine
                                Answer = RealtimeSearchEngine(f"{mode_context}\n\nQuery: {QueryFinal}", mode=mode)
                                break
                except Exception as e:
                    print(f"RealtimeSearchEngine error: {e}")
                    Answer = f"I encountered an issue while searching for information. Let me try a different approach."
            else:
                # Handle general queries with mode context
                # Status already set in process_text_message, but ensure it's visible
                if self.label.text() != "Thinking...":
                    self.label.setText("Thinking...")
                    self.update()
                    QApplication.processEvents()
                
                try:
                    # Check if user is asking "what can you do" BEFORE adding mode context
                    query_lower = message.lower().strip()
                    what_can_you_do_phrases = [
                        "what can you do", "what can you do for me", "what are your capabilities",
                        "what can you help with", "tell me what you can do", "what do you do",
                        "what are you capable of", "what services", "your capabilities", "your features"
                    ]
                    
                    # If asking "what can you do", pass original message directly with mode
                    # ChatBot will handle it and return early to prevent duplicates
                    if any(phrase in query_lower for phrase in what_can_you_do_phrases):
                        Answer = ChatBot(message, mode=mode)  # Pass original message, ChatBot handles it
                        # Don't process further - ChatBot already returned the response
                        if Answer:
                            return Answer  # Return immediately to prevent further processing
                    else:
                        # For other queries, add mode context
                        mode_context = f"[Mode: {mode}] {mode_prompt}\n\nUser Query: {message}"
                        
                        for queries in Decision:
                            if "general" in queries:
                                QueryFinal = queries.replace("general", "").strip()
                                # Use original message if QueryFinal is empty
                                query_to_use = QueryFinal if QueryFinal else message
                                # Pass mode to ChatBot for mode-specific behavior
                                Answer = ChatBot(f"{mode_context}\n\nQuery: {query_to_use}", mode=mode)
                                break
                            elif "exit" in queries:
                                Answer = "Goodbye! It was great working with you."
                                return Answer
                        
                        # If no general query found in Decision, try ChatBot with original message and mode
                        if not Answer:
                            Answer = ChatBot(message, mode=mode)
                except Exception as e:
                    print(f"ChatBot error: {e}")
                    # Fallback to RealtimeSearchEngine if ChatBot fails
                    try:
                        print("ChatBot failed, trying RealtimeSearchEngine...")
                        Answer = RealtimeSearchEngine(f"[Mode: {mode}] {message}")
                    except Exception as e2:
                        print(f"RealtimeSearchEngine fallback error: {e2}")
                        Answer = "I encountered an issue while processing your request. Let me try again."
            
            # Fallback if no answer was generated - try different approaches
            if not Answer or Answer.strip() == "":
                # Try ChatBot directly with mode
                try:
                    print("No answer from mode processing, trying ChatBot directly...")
                    Answer = ChatBot(message, mode=mode)
                except Exception as e:
                    print(f"ChatBot direct call error: {e}")
                    # Try RealtimeSearchEngine as last resort
                    try:
                        print("Trying RealtimeSearchEngine as fallback...")
                        Answer = RealtimeSearchEngine(message)
                    except Exception as e2:
                        print(f"RealtimeSearchEngine fallback error: {e2}")
                        Answer = "I'm here to help! Please try asking your question again or rephrase it. I can help with sales, real-time information, file analysis, and general questions."
            
            # Ensure answer mentions JARVIS, not FRIDAY
            if Answer and "FRIDAY" in Answer:
                Answer = Answer.replace("FRIDAY", "JARVIS")
                Answer = Answer.replace("friday", "JARVIS")
            
            return Answer
            
        except Exception as e:
            print(f"Error in process_with_mode: {e}")
            return f"I encountered an error while processing your request: {str(e)}"
    
    def pulse_status_label(self):
        """Create a pulsing animation effect for the status label"""
        try:
            if self.pulse_state == 0:
                self.label.setStyleSheet("""
                    color: #e0e0e0; 
                    font-size: 15px; 
                    font-weight: 600;
                    border: none; 
                    margin-top: -30px;
                    padding: 10px 20px;
                    background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                        stop:0 rgba(0, 220, 255, 0.25),
                        stop:1 rgba(120, 170, 255, 0.25));
                    border-radius: 25px;
                    border: 2px solid rgba(0, 220, 255, 0.6);
                    letter-spacing: 0.5px;
                """)
                self.pulse_state = 1
            else:
                self.label.setStyleSheet("""
                    color: #e0e0e0; 
                    font-size: 15px; 
                    font-weight: 600;
                    border: none; 
                    margin-top: -30px;
                    padding: 10px 20px;
                    background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                        stop:0 rgba(0, 200, 255, 0.15),
                        stop:1 rgba(100, 150, 255, 0.15));
                    border-radius: 25px;
                    border: 2px solid rgba(0, 200, 255, 0.4);
                    letter-spacing: 0.5px;
                """)
                self.pulse_state = 0
        except Exception as e:
            print(f"Error in pulse animation: {e}")



class InitialScreen(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        desktop = QApplication.desktop()
        screen_width = desktop.screenGeometry().width()
        screen_height = desktop.screenGeometry().height()
        content_layout = QVBoxLayout()
        content_layout.setContentsMargins(0, 0, 0, 0)

        gif_label = QLabel()
        gif_label.setAlignment(Qt.AlignCenter)
        gif_label.setStyleSheet("border: none; background: transparent;")
        gif_path = GraphicsDirPath + r'\FRIDAY.gif'
        if os.path.exists(gif_path):
            try:
                movie = QMovie(gif_path)
                if movie.isValid():
                    max_gif_size_H = int(screen_width / 16 * 9)
                    movie.setScaledSize(QSize(screen_width, max_gif_size_H))
                    movie.setCacheMode(QMovie.CacheAll)
                    gif_label.setMovie(movie)
                    gif_label.setFixedSize(screen_width, max_gif_size_H)
                    gif_label.setVisible(True)
                    movie.start()
                    print(f"Initial screen GIF loaded successfully from: {gif_path}")
                else:
                    print(f"Warning: Initial screen GIF file is not valid: {gif_path}")
                    gif_label.setFixedHeight(int(screen_width / 16 * 9))
            except Exception as e:
                print(f"Error loading initial screen GIF: {e}")
                gif_label.setText("")
                gif_label.setFixedHeight(int(screen_width / 16 * 9))
        else:
            gif_label.setText("")
            gif_label.setFixedHeight(int(screen_width / 16 * 9))
            print(f"Warning: GIF file not found at {gif_path}")

        gif_label.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.gif_label = gif_label  # Store reference
        self.icon_label = QLabel()
        pixmap = QPixmap(GraphicsDirPath + r'\Mic_on.png')  # Fixed this line
        new_pixmap = pixmap.scaled(60, 60)
        self.icon_label.setPixmap(new_pixmap)
        self.icon_label.setFixedSize(150, 150)
        self.icon_label.setAlignment(Qt.AlignCenter)
        self.toggled = True  # Start with mic ON for continuous listening
        # Initialize mic as ON
        MicButtonInitiated()
        self.icon_label.mousePressEvent = self.toggle_icon

        self.label = QLabel("")
        self.label.setStyleSheet("color: white; font-size: 16px; margin-bottom: 0;")
        content_layout.addWidget(gif_label, alignment=Qt.AlignCenter)
        content_layout.addWidget(self.label, alignment=Qt.AlignCenter)
        content_layout.addWidget(self.icon_label, alignment=Qt.AlignCenter)
        content_layout.setContentsMargins(0, 0, 0, 150)
        self.setLayout(content_layout)
        self.setFixedHeight(screen_height)
        self.setFixedWidth(screen_width)
        self.setStyleSheet("background-color: black;")
        
        self.timer = QTimer(self)
        self.timer.timeout.connect(self.SpeechRecogText)
        self.timer.start(5)

    def SpeechRecogText(self):
        with open(TempDirPath + r'\Status.data', 'r', encoding='utf-8') as file:  # Fixed this line
            messages = file.read()
            self.label.setText(messages)

    def load_icon(self, path, width=60, height=60):
        pixmap = QPixmap(path)
        new_pixmap = pixmap.scaled(width, height)
        self.icon_label.setPixmap(new_pixmap)

    def toggle_icon(self, event=None):
        if self.toggled:
            # Mic is currently ON, turning it OFF
            self.load_icon(GraphicsDirPath + r'\Mic_off.png', 60, 60)
            MicButtonClosed()
            self.toggled = False
        else:
            # Mic is currently OFF, turning it ON
            self.load_icon(GraphicsDirPath + r'\Mic_on.png', 60, 60)
            MicButtonInitiated()
            self.toggled = True


class MessageScreen(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        desktop = QApplication.desktop()
        screen_width = desktop.screenGeometry().width()
        screen_height = desktop.screenGeometry().height()
        layout = QVBoxLayout()
        label = QLabel("")
        layout.addWidget(label)
        chat_section = ChatSection()  # Ensure ChatSection is defined elsewhere
        # Set global reference so Main.py can add user messages
        from Frontend.GUI import set_global_chat_section
        set_global_chat_section(chat_section)
        layout.addWidget(chat_section)
        self.setLayout(layout)
        self.setStyleSheet("background-color: black;")
        self.setFixedHeight(screen_height)
        self.setFixedWidth(screen_width)


class CustomTopBar(QWidget):
    def __init__(self, parent, stacked_widget):
        super().__init__(parent)
        self.stacked_widget = stacked_widget
        self.current_screen = None
        self.initUI()

    def initUI(self):
        self.setFixedHeight(50)
        self.setStyleSheet("background-color: white;")
        layout = QHBoxLayout(self)
        layout.setContentsMargins(10, 5, 10, 5)
        layout.setSpacing(5)
        layout.setAlignment(Qt.AlignRight)

        home_button = QPushButton("Home")
        home_icon_path = GraphicsDirPath + r'\Home.png'
        if os.path.exists(home_icon_path):
            home_icon = QIcon(home_icon_path)
            home_button.setIcon(home_icon)
        home_button.setText("   Home")
        home_button.setStyleSheet("""
            QPushButton {
                height: 40px;
                line-height: 40px;
                background-color: white;
                color: black;
                border: 1px solid #ccc;
                border-radius: 4px;
                padding: 5px 15px;
                font-size: 14px;
                font-weight: 600;
            }
            QPushButton:hover {
                background-color: #f0f0f0;
            }
            QPushButton:pressed {
                background-color: #e0e0e0;
            }
        """)
        home_button.clicked.connect(self.showInitialScreen)

        message_button = QPushButton("Message")
        message_icon_path = GraphicsDirPath + r'\Message.png'
        if os.path.exists(message_icon_path):
            message_icon = QIcon(message_icon_path)
            message_button.setIcon(message_icon)
        message_button.setText("   Message")
        message_button.setStyleSheet("""
            QPushButton {
                height: 40px;
                line-height: 40px;
                background-color: white;
                color: black;
                border: 1px solid #ccc;
                border-radius: 4px;
                padding: 5px 15px;
                font-size: 14px;
                font-weight: 600;
            }
            QPushButton:hover {
                background-color: #f0f0f0;
            }
            QPushButton:pressed {
                background-color: #e0e0e0;
            }
        """)
        message_button.clicked.connect(self.showMessageScreen)

        minimize_button = QPushButton()
        minimize_icon_path = GraphicsDirPath + r'\Minimize.png'
        if os.path.exists(minimize_icon_path):
            minimize_icon = QIcon(minimize_icon_path)
            minimize_button.setIcon(minimize_icon)
        else:
            minimize_button.setText("−")
        minimize_button.setFixedSize(30, 30)
        minimize_button.setStyleSheet("""
            QPushButton {
                background-color: white;
                border: 1px solid #ccc;
                border-radius: 4px;
            }
            QPushButton:hover {
                background-color: #f0f0f0;
            }
            QPushButton:pressed {
                background-color: #e0e0e0;
            }
        """)
        minimize_button.clicked.connect(self.minimizeWindow)

        self.maximize_button = QPushButton()
        maximize_icon_path = GraphicsDirPath + r'\Maximize.png'
        restore_icon_path = GraphicsDirPath + r'\Restore.png'
        if os.path.exists(maximize_icon_path):
            self.maximize_icon = QIcon(maximize_icon_path)
            self.maximize_button.setIcon(self.maximize_icon)
        else:
            self.maximize_button.setText("□")
        if os.path.exists(restore_icon_path):
            self.restore_icon = QIcon(restore_icon_path)
        self.maximize_button.setFixedSize(30, 30)
        self.maximize_button.setStyleSheet("""
            QPushButton {
                background-color: white;
                border: 1px solid #ccc;
                border-radius: 4px;
            }
            QPushButton:hover {
                background-color: #f0f0f0;
            }
            QPushButton:pressed {
                background-color: #e0e0e0;
            }
        """)
        self.maximize_button.clicked.connect(self.maximizeWindow)

        close_button = QPushButton()
        close_icon_path = GraphicsDirPath + r'\Close.png'
        if os.path.exists(close_icon_path):
            close_icon = QIcon(close_icon_path)
            close_button.setIcon(close_icon)
        else:
            close_button.setText("×")
        close_button.setFixedSize(30, 30)
        close_button.setStyleSheet("""
            QPushButton {
                background-color: white;
                border: 1px solid #ccc;
                border-radius: 4px;
            }
            QPushButton:hover {
                background-color: #ff4444;
                color: white;
            }
            QPushButton:pressed {
                background-color: #cc0000;
            }
        """)
        close_button.clicked.connect(self.closeWindow)

        layout.addWidget(home_button)
        layout.addWidget(message_button)
        layout.addWidget(minimize_button)
        layout.addWidget(self.maximize_button)
        layout.addWidget(close_button)

        self.draggable = True
        self.offset = None

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.fillRect(self.rect(), Qt.white)
        super().paintEvent(event)



    def minimizeWindow(self):
        self.parent().showMinimized()

    def maximizeWindow(self):
        if self.parent().isMaximized():
            self.parent().showNormal()
            self.maximize_button.setIcon(self.maximize_icon)
        else:
            self.parent().showMaximized()
            self.maximize_button.setIcon(self.restore_icon)

    def closeWindow(self):
        self.parent().close()

    def showMessageScreen(self):
        self.stacked_widget.setCurrentIndex(1)

    def showInitialScreen(self):
        self.stacked_widget.setCurrentIndex(0)


class MainWindow(QMainWindow):
    def __init__(self):
        super(MainWindow, self).__init__()
        self.setWindowFlags(Qt.FramelessWindowHint | Qt.Window)
        self.initUI()

    def initUI(self):
        desktop = QApplication.desktop()
        screen_width = desktop.screenGeometry().width()
        screen_height = desktop.screenGeometry().height()
        stacked_widget = QStackedWidget(self)
        initial_screen = InitialScreen()
        message_screen = MessageScreen()
        stacked_widget.addWidget(initial_screen)
        stacked_widget.addWidget(message_screen)
        self.setGeometry(0, 0, screen_width, screen_height)
        self.setStyleSheet("background-color: black;")
        top_bar = CustomTopBar(self, stacked_widget)
        self.setMenuWidget(top_bar)
        self.setCentralWidget(stacked_widget)

# Global reference to chat section for adding messages from Main.py
_global_chat_section = None

def set_global_chat_section(chat_section):
    """Set global reference to chat section"""
    global _global_chat_section
    _global_chat_section = chat_section

def add_user_message_to_chat(message):
    """Add user message to chat from Main.py (voice input)"""
    global _global_chat_section
    
    if not message or not message.strip():
        print("add_user_message_to_chat: Empty message, skipping")
        return
    
    print(f"add_user_message_to_chat: Attempting to add message: '{message[:50]}...'")
    
    if _global_chat_section:
        try:
            print(f"add_user_message_to_chat: Global chat section found, adding message")
            _global_chat_section.addMessage(f"You: {message}", "White")
            # Force GUI update to ensure message is visible
            _global_chat_section.chat_text_edit.ensureCursorVisible()
            _global_chat_section.update()
            QApplication.processEvents()
            print(f"add_user_message_to_chat: Successfully added message to chat")
        except Exception as e:
            print(f"Error adding user message to chat: {e}")
            import traceback
            print(traceback.format_exc())
    else:
        print(f"WARNING: _global_chat_section is None - cannot add message to chat")

def clear_chat_history_global():
    """Clear chat history globally (can be called from Main.py)"""
    global _global_chat_section
    
    if _global_chat_section:
        try:
            _global_chat_section.clear_chat_history()
            print("Chat history cleared via global function")
        except Exception as e:
            print(f"Error clearing chat history: {e}")
            import traceback
            print(traceback.format_exc())
    else:
        print("WARNING: _global_chat_section is None - cannot clear chat history")
        # Try to find the chat section from the main window
        try:
            from PyQt5.QtWidgets import QApplication
            app = QApplication.instance()
            if app:
                for widget in app.allWidgets():
                    if hasattr(widget, 'chat_section') and widget.chat_section:
                        print(f"Found chat section via widget search, setting global reference")
                        set_global_chat_section(widget.chat_section)
                        widget.chat_section.clear_chat_history()
                        print(f"Successfully cleared chat history via widget search")
                        return
        except Exception as e2:
            print(f"Error in fallback chat section search: {e2}")

def GraphicalUserInterface():
    app = QApplication(sys.argv)
    window = MainWindow()
    # Set global chat section reference
    if hasattr(window, 'chat_section'):
        set_global_chat_section(window.chat_section)
    window.show()
    sys.exit(app.exec_())
# Run the application
if __name__ == "__main__":
    import sys
    app = QApplication(sys.argv)
    main_window = MainWindow()
    main_window.show()
    sys.exit(app.exec_())